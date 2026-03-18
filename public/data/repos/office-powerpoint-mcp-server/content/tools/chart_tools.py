"""
Chart data management tools for PowerPoint MCP Server.
Implements advanced chart data manipulation capabilities.
"""

from typing import Dict, List, Optional, Any
from mcp.types import ToolAnnotations
from pptx.chart.data import ChartData

def register_chart_tools(app, presentations, get_current_presentation_id, validate_parameters, 
                          is_positive, is_non_negative, is_in_range, is_valid_rgb):
    """Register chart data management tools with the FastMCP app."""
    
    @app.tool(
        annotations=ToolAnnotations(
            title="Update Chart Data",
        ),
    )
    def update_chart_data(
        slide_index: int,
        shape_index: int,
        categories: List[str],
        series_data: List[Dict],
        presentation_id: str = None
    ) -> Dict:
        """
        Replace existing chart data with new categories and series.
        
        Args:
            slide_index: Index of the slide (0-based)
            shape_index: Index of the chart shape (0-based)
            categories: List of category names
            series_data: List of dictionaries with 'name' and 'values' keys
            presentation_id: Optional presentation ID (uses current if not provided)
            
        Returns:
            Dictionary with operation results
        """
        try:
            # Get presentation
            pres_id = presentation_id or get_current_presentation_id()
            if pres_id not in presentations:
                return {"error": "Presentation not found"}
            
            pres = presentations[pres_id]
            
            # Validate slide index
            if not (0 <= slide_index < len(pres.slides)):
                return {"error": f"Slide index {slide_index} out of range"}
            
            slide = pres.slides[slide_index]
            
            # Validate shape index
            if not (0 <= shape_index < len(slide.shapes)):
                return {"error": f"Shape index {shape_index} out of range"}
            
            shape = slide.shapes[shape_index]
            
            # Check if shape is a chart
            if not hasattr(shape, 'has_chart') or not shape.has_chart:
                return {"error": "Shape is not a chart"}
            
            chart = shape.chart
            
            # Create new ChartData
            chart_data = ChartData()
            chart_data.categories = categories
            
            # Add series data
            for series in series_data:
                if 'name' not in series or 'values' not in series:
                    return {"error": "Each series must have 'name' and 'values' keys"}
                
                chart_data.add_series(series['name'], series['values'])
            
            # Replace chart data
            chart.replace_data(chart_data)
            
            return {
                "message": f"Updated chart data on slide {slide_index}, shape {shape_index}",
                "categories": categories,
                "series_count": len(series_data),
                "series_names": [s['name'] for s in series_data]
            }
            
        except Exception as e:
            return {"error": f"Failed to update chart data: {str(e)}"}