"""
Enhanced template management utilities for PowerPoint MCP Server.
Advanced slide creation with dynamic sizing, auto-wrapping, and visual effects.
Combines features from both basic and enhanced template systems.
"""
import json
import os
import re
from typing import Dict, List, Optional, Any, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import utils.content_utils as content_utils
import utils.design_utils as design_utils


class TextSizeCalculator:
    """Calculate optimal text sizes based on content and container dimensions."""
    
    def __init__(self):
        self.character_widths = {
            'narrow': 0.6,  # i, l, t
            'normal': 1.0,  # most characters
            'wide': 1.3,    # m, w
            'space': 0.5    # space character
        }
    
    def estimate_text_width(self, text: str, font_size: int) -> float:
        """Estimate text width in points based on character analysis."""
        if not text:
            return 0
        
        width = 0
        for char in text:
            if char in 'iltj':
                width += self.character_widths['narrow']
            elif char in 'mwMW':
                width += self.character_widths['wide']
            elif char == ' ':
                width += self.character_widths['space']
            else:
                width += self.character_widths['normal']
        
        return width * font_size * 0.6  # Approximation factor
    
    def estimate_text_height(self, text: str, font_size: int, line_spacing: float = 1.2) -> float:
        """Estimate text height based on line count and spacing."""
        lines = len(text.split('\n'))
        return lines * font_size * line_spacing * 1.3  # Convert to points
    
    def calculate_optimal_font_size(self, text: str, container_width: float, 
                                  container_height: float, font_type: str = 'body',
                                  min_size: int = 8, max_size: int = 36) -> int:
        """Calculate optimal font size to fit text in container."""
        container_width_pts = container_width * 72  # Convert inches to points
        container_height_pts = container_height * 72
        
        # Start with a reasonable size and adjust
        for font_size in range(max_size, min_size - 1, -1):
            estimated_width = self.estimate_text_width(text, font_size)
            estimated_height = self.estimate_text_height(text, font_size)
            
            if estimated_width <= container_width_pts * 0.9 and estimated_height <= container_height_pts * 0.9:
                return font_size
        
        return min_size
    
    def wrap_text_intelligently(self, text: str, max_width: float, font_size: int) -> str:
        """Intelligently wrap text to fit within specified width."""
        if not text:
            return text
        
        max_width_pts = max_width * 72
        words = text.split()
        wrapped_lines = []
        current_line = []
        
        for word in words:
            test_line = current_line + [word]
            test_text = ' '.join(test_line)
            
            if self.estimate_text_width(test_text, font_size) <= max_width_pts:
                current_line.append(word)
            else:
                if current_line:
                    wrapped_lines.append(' '.join(current_line))
                    current_line = [word]
                else:
                    # Single word is too long, force wrap
                    wrapped_lines.append(word)
        
        if current_line:
            wrapped_lines.append(' '.join(current_line))
        
        return '\n'.join(wrapped_lines)


class VisualEffectsManager:
    """Manage and apply visual effects to PowerPoint elements."""
    
    def __init__(self, templates_data: Dict):
        self.templates_data = templates_data
        self.text_effects = templates_data.get('text_effects', {})
        self.image_effects = templates_data.get('image_effects', {})
    
    def apply_text_effects(self, text_frame, effects: List[str], color_scheme: str) -> None:
        """Apply text effects like shadows, glows, and outlines."""
        for effect_name in effects:
            if effect_name not in self.text_effects:
                continue
            
            effect_config = self.text_effects[effect_name]
            effect_type = effect_config.get('type')
            
            # Note: These are simplified implementations
            # Full implementation would require XML manipulation
            try:
                if effect_type == 'shadow':
                    self._apply_text_shadow(text_frame, effect_config, color_scheme)
                elif effect_type == 'glow':
                    self._apply_text_glow(text_frame, effect_config, color_scheme)
                elif effect_type == 'outline':
                    self._apply_text_outline(text_frame, effect_config, color_scheme)
            except Exception:
                # Graceful fallback if effect application fails
                pass
    
    def _apply_text_shadow(self, text_frame, config: Dict, color_scheme: str) -> None:
        """Apply shadow effect to text (simplified implementation)."""
        # In a full implementation, this would manipulate the XML directly
        # For now, we'll apply basic formatting that creates a shadow-like effect
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                # Make text slightly bolder to simulate shadow depth
                run.font.bold = True
    
    def _apply_text_glow(self, text_frame, config: Dict, color_scheme: str) -> None:
        """Apply glow effect to text (simplified implementation)."""
        pass  # Would require XML manipulation for true glow effect
    
    def _apply_text_outline(self, text_frame, config: Dict, color_scheme: str) -> None:
        """Apply outline effect to text (simplified implementation)."""
        pass  # Would require XML manipulation for true outline effect
    
    def apply_image_effects(self, image_shape, effect_name: str, color_scheme: str) -> None:
        """Apply visual effects to image shapes."""
        if effect_name not in self.image_effects:
            return
        
        effect_config = self.image_effects[effect_name]
        
        try:
            # Apply shadow if specified
            if 'shadow' in effect_config:
                shadow_config = effect_config['shadow']
                # Simplified shadow application
                pass
            
            # Apply border if specified
            if 'border' in effect_config:
                border_config = effect_config['border']
                if 'width' in border_config:
                    image_shape.line.width = Pt(border_config['width'])
                if 'color_role' in border_config:
                    color = self._get_color_from_scheme(color_scheme, border_config['color_role'])
                    image_shape.line.color.rgb = RGBColor(*color)
                elif 'color' in border_config:
                    image_shape.line.color.rgb = RGBColor(*border_config['color'])
        
        except Exception:
            # Graceful fallback
            pass
    
    def _get_color_from_scheme(self, color_scheme: str, color_role: str) -> Tuple[int, int, int]:
        """Get color from scheme (helper method)."""
        schemes = self.templates_data.get('color_schemes', {})
        if color_scheme in schemes and color_role in schemes[color_scheme]:
            return tuple(schemes[color_scheme][color_role])
        return (0, 0, 0)  # Default black


class EnhancedTemplateManager:
    """Enhanced template manager with dynamic features."""
    
    def __init__(self, template_file_path: str = None):
        self.text_calculator = TextSizeCalculator()
        self.load_templates(template_file_path)
        self.effects_manager = VisualEffectsManager(self.templates_data)
    
    def load_templates(self, template_file_path: str = None) -> None:
        """Load unified templates with all dynamic features."""
        if template_file_path is None:
            current_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            # Use the unified template file
            template_file_path = os.path.join(current_dir, 'slide_layout_templates.json')
        
        try:
            with open(template_file_path, 'r', encoding='utf-8') as f:
                self.templates_data = json.load(f)
        except FileNotFoundError:
            raise FileNotFoundError(f"Template file not found: {template_file_path}")
        except json.JSONDecodeError as e:
            raise ValueError(f"Invalid JSON in template file: {str(e)}")


    def get_dynamic_font_size(self, element: Dict, content: str = None) -> int:
        """Calculate dynamic font size based on content and container."""
        content = content or element.get('placeholder_text', '')
        if not content:
            return 14  # Default size
        
        # Get container dimensions
        pos = element.get('position', {})
        container_width = pos.get('width', 4.0)
        container_height = pos.get('height', 1.0)
        
        # Get font constraints
        font_type = element.get('styling', {}).get('font_type', 'body')
        sizing_rules = self.templates_data.get('auto_sizing_rules', {})
        base_sizes = sizing_rules.get('text_measurement', {}).get('base_font_sizes', {})
        
        if font_type in base_sizes:
            min_size = base_sizes[font_type]['min']
            max_size = base_sizes[font_type]['max']
            default_size = base_sizes[font_type]['default']
        else:
            min_size, max_size, default_size = 10, 18, 14
        
        # Check if dynamic sizing is requested
        font_size_setting = element.get('styling', {}).get('font_size')
        if font_size_setting == 'dynamic':
            return self.text_calculator.calculate_optimal_font_size(
                content, container_width, container_height, font_type, min_size, max_size
            )
        
        return default_size
    
    def apply_enhanced_slide_template(self, slide, template_id: str, color_scheme: str = 'modern_blue',
                                    content_mapping: Dict = None, image_paths: Dict = None) -> Dict:
        """Apply enhanced slide template with all dynamic features."""
        try:
            if template_id not in self.templates_data.get('templates', {}):
                # Fall back to regular template application
                return apply_slide_template_basic(slide, template_id, color_scheme, content_mapping, image_paths)
            
            template = self.templates_data['templates'][template_id]
            elements_created = []
            
            # Apply enhanced background if specified
            background_config = template.get('background')
            if background_config:
                apply_slide_background(slide, background_config, self.templates_data, color_scheme)
            
            # Create enhanced elements
            for element in template.get('elements', []):
                element_type = element.get('type')
                element_role = element.get('role', '')
                
                try:
                    # Override content if provided
                    custom_content = None
                    if content_mapping and element_role in content_mapping:
                        custom_content = content_mapping[element_role]
                    
                    created_element = None
                    
                    if element_type == 'text':
                        created_element = self.create_enhanced_text_element(
                            slide, element, self.templates_data, color_scheme, custom_content
                        )
                    elif element_type == 'shape':
                        created_element = create_shape_element(slide, element, self.templates_data, color_scheme)
                    elif element_type == 'image':
                        image_path = image_paths.get(element_role) if image_paths else None
                        created_element = create_image_element(slide, element, image_path)
                    elif element_type == 'table':
                        created_element = create_table_element(slide, element, self.templates_data, color_scheme)
                    elif element_type == 'chart':
                        created_element = create_chart_element(slide, element, self.templates_data, color_scheme)
                    
                    if created_element:
                        elements_created.append({
                            'type': element_type,
                            'role': element_role,
                            'index': len(slide.shapes) - 1,
                            'enhanced_features': self.get_element_features(element)
                        })
                
                except Exception as e:
                    elements_created.append({
                        'type': element_type,
                        'role': element_role,
                        'error': str(e)
                    })
            
            return {
                'success': True,
                'template_id': template_id,
                'template_name': template.get('name', template_id),
                'color_scheme': color_scheme,
                'elements_created': elements_created,
                'enhanced_features_applied': [
                    'Dynamic text sizing',
                    'Automatic text wrapping',
                    'Visual effects',
                    'Intelligent content adaptation'
                ]
            }
        
        except Exception as e:
            return {
                'success': False,
                'error': f"Failed to apply enhanced template: {str(e)}"
            }
    
    def create_enhanced_text_element(self, slide, element: Dict, templates_data: Dict, 
                                   color_scheme: str, custom_content: str = None) -> Any:
        """Create text element with enhanced features."""
        pos = element['position']
        
        # Determine content
        content = custom_content or element.get('placeholder_text', '')
        
        # Apply auto-wrapping if enabled
        styling = element.get('styling', {})
        if styling.get('auto_wrap', False):
            container_width = pos.get('width', 4.0)
            font_size = self.get_dynamic_font_size(element, content)
            content = self.text_calculator.wrap_text_intelligently(content, container_width, font_size)
        
        # Create text box
        textbox = slide.shapes.add_textbox(
            Inches(pos['left']), 
            Inches(pos['top']), 
            Inches(pos['width']), 
            Inches(pos['height'])
        )
        
        textbox.text_frame.text = content
        textbox.text_frame.word_wrap = True
        
        # Apply dynamic font sizing
        font_size = self.get_dynamic_font_size(element, content)
        
        # Apply enhanced styling
        self.apply_enhanced_text_styling(textbox.text_frame, element, templates_data, color_scheme, font_size)
        
        # Apply auto-fit if enabled
        if styling.get('auto_fit', False):
            textbox.text_frame.auto_size = True
        
        return textbox
    
    def apply_enhanced_text_styling(self, text_frame, element: Dict, templates_data: Dict, 
                                  color_scheme: str, font_size: int) -> None:
        """Apply enhanced text styling with effects and dynamic features."""
        styling = element.get('styling', {})
        
        # Get typography style
        typography_style = templates_data.get('typography_styles', {}).get('modern_sans', {})
        font_type = styling.get('font_type', 'body')
        font_config = typography_style.get(font_type, {'name': 'Segoe UI', 'weight': 'normal'})
        
        # Color handling
        color = None
        if 'color_role' in styling:
            color = get_color_from_scheme(templates_data, color_scheme, styling['color_role'])
        elif 'color' in styling:
            color = tuple(styling['color'])
        
        # Alignment mapping
        alignment_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        
        # Vertical alignment mapping
        vertical_alignment_map = {
            'top': MSO_VERTICAL_ANCHOR.TOP,
            'middle': MSO_VERTICAL_ANCHOR.MIDDLE,
            'bottom': MSO_VERTICAL_ANCHOR.BOTTOM
        }
        
        # Apply vertical alignment to text frame
        if 'vertical_alignment' in styling:
            v_align = styling['vertical_alignment']
            if v_align in vertical_alignment_map:
                text_frame.vertical_anchor = vertical_alignment_map[v_align]
        
        # Dynamic line spacing
        line_spacing = styling.get('line_spacing', 1.2)
        if line_spacing == 'dynamic':
            content_length = len(text_frame.text)
            if content_length > 300:
                line_spacing = 1.4
            elif content_length > 150:
                line_spacing = 1.3
            else:
                line_spacing = 1.2
        
        # Apply formatting to paragraphs and runs
        for paragraph in text_frame.paragraphs:
            # Set alignment
            if 'alignment' in styling and styling['alignment'] in alignment_map:
                paragraph.alignment = alignment_map[styling['alignment']]
            
            # Set line spacing
            paragraph.line_spacing = line_spacing
            
            # Apply formatting to runs
            for run in paragraph.runs:
                font = run.font
                
                # Font family and size
                font.name = font_config['name']
                font.size = Pt(font_size)
                
                # Font weight and style
                weight = font_config.get('weight', 'normal')
                font.bold = styling.get('bold', weight in ['bold', 'semibold'])
                font.italic = styling.get('italic', font_config.get('style') == 'italic')
                font.underline = styling.get('underline', False)
                
                # Color
                if color:
                    font.color.rgb = RGBColor(*color)
        
        # Apply text effects
        text_effects = styling.get('text_effects', [])
        if text_effects:
            self.effects_manager.apply_text_effects(text_frame, text_effects, color_scheme)
    
    def get_element_features(self, element: Dict) -> List[str]:
        """Get list of enhanced features applied to an element."""
        features = []
        styling = element.get('styling', {})
        
        if styling.get('font_size') == 'dynamic':
            features.append('Dynamic text sizing')
        if styling.get('auto_wrap'):
            features.append('Automatic text wrapping')
        if styling.get('text_effects'):
            features.append('Text visual effects')
        if styling.get('auto_fit'):
            features.append('Auto-fit content')
        if 'fill_gradient' in styling:
            features.append('Gradient fills')
        if styling.get('shadow') or styling.get('glow'):
            features.append('Advanced visual effects')
        
        return features


# Global instance for enhanced features
enhanced_template_manager = EnhancedTemplateManager()


def get_enhanced_template_manager() -> EnhancedTemplateManager:
    """Get the global enhanced template manager instance."""
    return enhanced_template_manager


def calculate_dynamic_font_size(text: str, container_width: float, container_height: float, 
                               font_type: str = 'body') -> int:
    """Calculate optimal font size for given text and container."""
    return enhanced_template_manager.text_calculator.calculate_optimal_font_size(
        text, container_width, container_height, font_type
    )


def wrap_text_automatically(text: str, container_width: float, font_size: int) -> str:
    """Automatically wrap text to fit container width."""
    return enhanced_template_manager.text_calculator.wrap_text_intelligently(
        text, container_width, font_size
    )


def load_slide_templates(template_file_path: str = None) -> Dict:
    """
    Load slide layout templates from JSON file.
    
    Args:
        template_file_path: Path to template JSON file (defaults to slide_layout_templates.json)
        
    Returns:
        Dictionary containing all template definitions
    """
    if template_file_path is None:
        # Default to the template file in the same directory as the script
        current_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        template_file_path = os.path.join(current_dir, 'slide_layout_templates.json')
    
    try:
        with open(template_file_path, 'r', encoding='utf-8') as f:
            templates = json.load(f)
        return templates
    except FileNotFoundError:
        raise FileNotFoundError(f"Template file not found: {template_file_path}")
    except json.JSONDecodeError as e:
        raise ValueError(f"Invalid JSON in template file: {str(e)}")


def get_available_templates() -> List[Dict]:
    """
    Get a list of all available slide templates.
    
    Returns:
        List of template information dictionaries
    """
    try:
        templates_data = load_slide_templates()
        template_list = []
        
        for template_id, template_info in templates_data.get('templates', {}).items():
            template_list.append({
                'id': template_id,
                'name': template_info.get('name', template_id),
                'description': template_info.get('description', ''),
                'layout_type': template_info.get('layout_type', 'content'),
                'element_count': len(template_info.get('elements', []))
            })
        
        return template_list
    except Exception as e:
        return [{'error': f"Failed to load templates: {str(e)}"}]


def get_color_from_scheme(templates_data: Dict, color_scheme: str, color_role: str) -> Tuple[int, int, int]:
    """
    Get RGB color values from a color scheme.
    
    Args:
        templates_data: Template data dictionary
        color_scheme: Name of the color scheme
        color_role: Role of the color (primary, secondary, accent1, etc.)
        
    Returns:
        RGB color tuple (r, g, b)
    """
    color_schemes = templates_data.get('color_schemes', {})
    
    if color_scheme not in color_schemes:
        color_scheme = 'modern_blue'  # Default fallback
    
    scheme = color_schemes[color_scheme]
    return tuple(scheme.get(color_role, scheme.get('primary', [0, 120, 215])))


def get_font_settings(templates_data: Dict, font_type: str, font_size: str) -> Dict:
    """
    Get font settings from typography configuration.
    
    Args:
        templates_data: Template data dictionary
        font_type: Type of font (title, subtitle, body, caption)
        font_size: Size category (large, medium, small)
        
    Returns:
        Dictionary with font settings
    """
    typography = templates_data.get('typography', {})
    
    if font_type not in typography:
        font_type = 'body'  # Default fallback
    
    font_config = typography[font_type]
    size_key = f'font_size_{font_size}'
    
    return {
        'name': font_config.get('font_name', 'Segoe UI'),
        'size': font_config.get(size_key, font_config.get('font_size_medium', 14)),
        'bold': font_config.get('bold', False)
    }


def apply_text_styling(text_frame, styling: Dict, templates_data: Dict, color_scheme: str) -> None:
    """
    Apply text styling based on template configuration.
    
    Args:
        text_frame: PowerPoint text frame object
        styling: Styling configuration from template
        templates_data: Template data dictionary
        color_scheme: Selected color scheme
    """
    # Get font settings
    font_type = styling.get('font_type', 'body')
    font_size_category = styling.get('font_size', 'medium')
    font_settings = get_font_settings(templates_data, font_type, font_size_category)
    
    # Get color
    color = None
    if 'color_role' in styling:
        color = get_color_from_scheme(templates_data, color_scheme, styling['color_role'])
    elif 'color' in styling:
        color = tuple(styling['color'])
    
    # Apply alignment
    alignment_map = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
        'justify': PP_ALIGN.JUSTIFY
    }
    
    # Apply formatting to all paragraphs and runs
    for paragraph in text_frame.paragraphs:
        if 'alignment' in styling and styling['alignment'] in alignment_map:
            paragraph.alignment = alignment_map[styling['alignment']]
        
        for run in paragraph.runs:
            font = run.font
            font.name = font_settings['name']
            font.size = Pt(font_settings['size'])
            font.bold = styling.get('bold', font_settings['bold'])
            font.italic = styling.get('italic', False)
            font.underline = styling.get('underline', False)
            
            if color:
                font.color.rgb = RGBColor(*color)


def create_text_element(slide, element: Dict, templates_data: Dict, color_scheme: str) -> Any:
    """
    Create a text element on a slide based on template configuration.
    
    Args:
        slide: PowerPoint slide object
        element: Element configuration from template
        templates_data: Template data dictionary
        color_scheme: Selected color scheme
        
    Returns:
        Created text box shape
    """
    pos = element['position']
    textbox = slide.shapes.add_textbox(
        Inches(pos['left']), 
        Inches(pos['top']), 
        Inches(pos['width']), 
        Inches(pos['height'])
    )
    
    # Set text content
    textbox.text_frame.text = element.get('placeholder_text', '')
    
    # Apply styling
    styling = element.get('styling', {})
    apply_text_styling(textbox.text_frame, styling, templates_data, color_scheme)
    
    return textbox


def create_image_element(slide, element: Dict, image_path: str = None) -> Any:
    """
    Create an image element on a slide based on template configuration.
    
    Args:
        slide: PowerPoint slide object
        element: Element configuration from template
        image_path: Optional path to image file
        
    Returns:
        Created image shape or None if no image provided
    """
    if not image_path:
        # Create placeholder rectangle if no image provided
        pos = element['position']
        placeholder = slide.shapes.add_shape(
            1,  # Rectangle shape
            Inches(pos['left']),
            Inches(pos['top']),
            Inches(pos['width']),
            Inches(pos['height'])
        )
        
        # Add placeholder text
        if hasattr(placeholder, 'text_frame'):
            placeholder.text_frame.text = element.get('placeholder_text', 'Image Placeholder')
        
        return placeholder
    
    try:
        pos = element['position']
        image_shape = content_utils.add_image(
            slide, 
            image_path, 
            pos['left'], 
            pos['top'], 
            pos['width'], 
            pos['height']
        )
        
        # Apply styling if specified
        styling = element.get('styling', {})
        if styling.get('shadow'):
            # Apply shadow effect (simplified)
            pass
        
        return image_shape
    except Exception:
        # Fallback to placeholder if image fails to load
        return create_image_element(slide, element, None)


def create_shape_element(slide, element: Dict, templates_data: Dict, color_scheme: str) -> Any:
    """
    Create a shape element on a slide based on template configuration.
    
    Args:
        slide: PowerPoint slide object
        element: Element configuration from template
        templates_data: Template data dictionary
        color_scheme: Selected color scheme
        
    Returns:
        Created shape
    """
    pos = element['position']
    shape_type = element.get('shape_type', 'rectangle')
    
    try:
        # Import the shape creation function from the main server
        from ppt_mcp_server import add_shape_direct
        shape = add_shape_direct(slide, shape_type, pos['left'], pos['top'], pos['width'], pos['height'])
        
        # Apply styling
        styling = element.get('styling', {})
        
        # Fill color
        if 'fill_color_role' in styling:
            fill_color = get_color_from_scheme(templates_data, color_scheme, styling['fill_color_role'])
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*fill_color)
        elif 'fill_color' in styling:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*styling['fill_color'])
        
        # Line color
        if 'line_color_role' in styling:
            line_color = get_color_from_scheme(templates_data, color_scheme, styling['line_color_role'])
            shape.line.color.rgb = RGBColor(*line_color)
        elif styling.get('no_border'):
            shape.line.fill.background()
        
        # Transparency
        if 'transparency' in styling:
            # Note: Transparency implementation would need additional XML manipulation
            pass
        
        return shape
    except Exception as e:
        # Create a simple rectangle as fallback
        textbox = slide.shapes.add_textbox(
            Inches(pos['left']), 
            Inches(pos['top']), 
            Inches(pos['width']), 
            Inches(pos['height'])
        )
        textbox.text_frame.text = f"Shape: {shape_type}"
        return textbox


def create_table_element(slide, element: Dict, templates_data: Dict, color_scheme: str) -> Any:
    """
    Create a table element on a slide based on template configuration.
    
    Args:
        slide: PowerPoint slide object
        element: Element configuration from template
        templates_data: Template data dictionary
        color_scheme: Selected color scheme
        
    Returns:
        Created table shape
    """
    pos = element['position']
    table_config = element.get('table_config', {})
    
    rows = table_config.get('rows', 3)
    cols = table_config.get('cols', 3)
    
    # Create table
    table_shape = content_utils.add_table(
        slide, rows, cols, pos['left'], pos['top'], pos['width'], pos['height']
    )
    table = table_shape.table
    
    # Populate with data if provided
    data = table_config.get('data', [])
    for r in range(min(rows, len(data))):
        for c in range(min(cols, len(data[r]))):
            table.cell(r, c).text = str(data[r][c])
    
    # Apply styling
    styling = element.get('styling', {})
    header_row = table_config.get('header_row', True)
    
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            
            if r == 0 and header_row:
                # Header styling
                if 'header_bg_color_role' in styling:
                    bg_color = get_color_from_scheme(templates_data, color_scheme, styling['header_bg_color_role'])
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(*bg_color)
                
                # Header text color
                if 'header_text_color' in styling:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(*styling['header_text_color'])
                            run.font.bold = True
            else:
                # Body styling
                if 'body_bg_color_role' in styling:
                    bg_color = get_color_from_scheme(templates_data, color_scheme, styling['body_bg_color_role'])
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(*bg_color)
    
    return table_shape


def create_chart_element(slide, element: Dict, templates_data: Dict, color_scheme: str) -> Any:
    """
    Create a chart element on a slide based on template configuration.
    
    Args:
        slide: PowerPoint slide object
        element: Element configuration from template
        templates_data: Template data dictionary
        color_scheme: Selected color scheme
        
    Returns:
        Created chart object
    """
    pos = element['position']
    chart_config = element.get('chart_config', {})
    
    chart_type = chart_config.get('type', 'column')
    categories = chart_config.get('categories', ['A', 'B', 'C'])
    series_data = chart_config.get('series', [{'name': 'Series 1', 'values': [1, 2, 3]}])
    
    # Extract series names and values
    series_names = [s['name'] for s in series_data]
    series_values = [s['values'] for s in series_data]
    
    try:
        # Create chart
        chart = content_utils.add_chart(
            slide, chart_type, pos['left'], pos['top'], pos['width'], pos['height'],
            categories, series_names, series_values
        )
        
        # Apply formatting
        chart_title = chart_config.get('title')
        if chart_title:
            content_utils.format_chart(chart, title=chart_title)
        
        return chart
    except Exception as e:
        # Create placeholder if chart creation fails
        textbox = slide.shapes.add_textbox(
            Inches(pos['left']), 
            Inches(pos['top']), 
            Inches(pos['width']), 
            Inches(pos['height'])
        )
        textbox.text_frame.text = f"Chart: {chart_type}\n{chart_title or 'Chart Placeholder'}"
        return textbox


def apply_slide_background(slide, background_config: Dict, templates_data: Dict, color_scheme: str) -> None:
    """
    Apply background styling to a slide based on template configuration.
    
    Args:
        slide: PowerPoint slide object
        background_config: Background configuration from template
        templates_data: Template data dictionary
        color_scheme: Selected color scheme
    """
    if not background_config:
        return
    
    bg_type = background_config.get('type', 'solid')
    
    if bg_type == 'professional_gradient':
        style = background_config.get('style', 'subtle')
        direction = background_config.get('direction', 'diagonal')
        design_utils.create_professional_gradient_background(slide, color_scheme, style, direction)
    elif bg_type == 'solid':
        color_role = background_config.get('color_role', 'light')
        # Note: Solid background would require XML manipulation for proper implementation
        pass




def apply_slide_template_basic(slide, template_id: str, color_scheme: str = 'modern_blue', 
                              content_mapping: Dict = None, image_paths: Dict = None) -> Dict:
    """
    Apply a basic slide template to create a formatted slide.
    
    Args:
        slide: PowerPoint slide object
        template_id: ID of the template to apply
        color_scheme: Color scheme to use
        content_mapping: Dictionary mapping element roles to content
        image_paths: Dictionary mapping image element roles to file paths
        
    Returns:
        Dictionary with application results
    """
    try:
        # Load templates
        templates_data = load_slide_templates()
        
        if template_id not in templates_data.get('templates', {}):
            return {
                'success': False,
                'error': f"Template '{template_id}' not found"
            }
        
        template = templates_data['templates'][template_id]
        elements_created = []
        
        # Apply background if specified
        background_config = template.get('background')
        if background_config:
            apply_slide_background(slide, background_config, templates_data, color_scheme)
        
        # Create elements
        for element in template.get('elements', []):
            element_type = element.get('type')
            element_role = element.get('role', '')
            
            try:
                # Override placeholder text with custom content if provided
                if content_mapping and element_role in content_mapping:
                    element = element.copy()  # Don't modify original template
                    element['placeholder_text'] = content_mapping[element_role]
                
                created_element = None
                
                if element_type == 'text':
                    created_element = create_text_element(slide, element, templates_data, color_scheme)
                elif element_type == 'image':
                    image_path = image_paths.get(element_role) if image_paths else None
                    created_element = create_image_element(slide, element, image_path)
                elif element_type == 'shape':
                    created_element = create_shape_element(slide, element, templates_data, color_scheme)
                elif element_type == 'table':
                    created_element = create_table_element(slide, element, templates_data, color_scheme)
                elif element_type == 'chart':
                    created_element = create_chart_element(slide, element, templates_data, color_scheme)
                
                if created_element:
                    elements_created.append({
                        'type': element_type,
                        'role': element_role,
                        'index': len(slide.shapes) - 1
                    })
                    
            except Exception as e:
                # Continue with other elements if one fails
                elements_created.append({
                    'type': element_type,
                    'role': element_role,
                    'error': str(e)
                })
        
        return {
            'success': True,
            'template_id': template_id,
            'template_name': template.get('name', template_id),
            'color_scheme': color_scheme,
            'elements_created': elements_created,
            'total_elements': len(template.get('elements', []))
        }
        
    except Exception as e:
        return {
            'success': False,
            'error': f"Failed to apply template: {str(e)}"
        }


def apply_slide_template(slide, template_id: str, color_scheme: str = 'modern_blue', 
                        content_mapping: Dict = None, image_paths: Dict = None) -> Dict:
    """
    Apply a slide template with all enhanced features.
    
    Args:
        slide: PowerPoint slide object
        template_id: ID of the template to apply
        color_scheme: Color scheme to use
        content_mapping: Dictionary mapping element roles to content
        image_paths: Dictionary mapping image element roles to file paths
        
    Returns:
        Dictionary with application results
    """
    # All templates now have enhanced features built-in
    return enhanced_template_manager.apply_enhanced_slide_template(
        slide, template_id, color_scheme, content_mapping, image_paths
    )


def create_presentation_from_template_sequence(presentation: Presentation, template_sequence: List[Dict],
                                              color_scheme: str = 'modern_blue') -> Dict:
    """
    Create a complete presentation from a sequence of templates.
    
    Args:
        presentation: PowerPoint presentation object
        template_sequence: List of template configurations
        color_scheme: Color scheme to apply to all slides
        
    Returns:
        Dictionary with creation results
    """
    results = {
        'success': True,
        'slides_created': [],
        'total_slides': len(template_sequence),
        'color_scheme': color_scheme
    }
    
    for i, slide_config in enumerate(template_sequence):
        try:
            # Get template configuration
            template_id = slide_config.get('template_id')
            content_mapping = slide_config.get('content', {})
            image_paths = slide_config.get('images', {})
            
            if not template_id:
                results['slides_created'].append({
                    'slide_index': i,
                    'success': False,
                    'error': 'No template_id specified'
                })
                continue
            
            # Add new slide (using layout 1 as default content layout)
            layout = presentation.slide_layouts[1]
            slide = presentation.slides.add_slide(layout)
            
            # Apply template
            template_result = apply_slide_template(
                slide, template_id, color_scheme, content_mapping, image_paths
            )
            
            template_result['slide_index'] = i
            results['slides_created'].append(template_result)
            
            if not template_result['success']:
                results['success'] = False
                
        except Exception as e:
            results['slides_created'].append({
                'slide_index': i,
                'success': False,
                'error': f"Failed to create slide {i}: {str(e)}"
            })
            results['success'] = False
    
    return results


def get_template_usage_examples() -> Dict:
    """
    Get examples of how to use different templates.
    
    Returns:
        Dictionary with usage examples
    """
    return {
        "single_slide_example": {
            "description": "Apply a single template to a slide",
            "code": {
                "template_id": "text_with_image",
                "color_scheme": "modern_blue",
                "content_mapping": {
                    "title": "Our Solution",
                    "content": "• Increased efficiency by 40%\n• Reduced costs significantly\n• Improved user satisfaction",
                },
                "image_paths": {
                    "supporting": "/path/to/solution_image.jpg"
                }
            }
        },
        "presentation_sequence_example": {
            "description": "Create a complete presentation from templates",
            "code": [
                {
                    "template_id": "title_slide",
                    "content": {
                        "title": "2024 Business Review",
                        "subtitle": "Annual Performance Report",
                        "author": "John Smith, CEO"
                    }
                },
                {
                    "template_id": "agenda_slide",
                    "content": {
                        "agenda_items": "1. Executive Summary\n\n2. Financial Performance\n\n3. Market Analysis\n\n4. Future Strategy"
                    }
                },
                {
                    "template_id": "key_metrics_dashboard",
                    "content": {
                        "metric_1_value": "92%",
                        "metric_2_value": "$3.2M",
                        "metric_3_value": "340",
                        "metric_4_value": "18%"
                    }
                },
                {
                    "template_id": "thank_you_slide",
                    "content": {
                        "contact": "Questions?\njohn.smith@company.com\n(555) 123-4567"
                    }
                }
            ]
        },
        "available_templates": [
            "title_slide", "text_with_image", "two_column_text", "two_column_text_images",
            "three_column_layout", "agenda_slide", "chapter_intro", "thank_you_slide",
            "timeline_slide", "data_table_slide", "chart_comparison", "full_image_slide",
            "process_flow", "quote_testimonial", "key_metrics_dashboard", 
            "before_after_comparison", "team_introduction"
        ],
        "color_schemes": [
            "modern_blue", "corporate_gray", "elegant_green", "warm_red"
        ]
    }