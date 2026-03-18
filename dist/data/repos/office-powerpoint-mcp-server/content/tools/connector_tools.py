"""
Connector and line tools for PowerPoint MCP Server.
Implements connector line/arrow drawing capabilities.
"""

from typing import Dict, List, Optional, Any
from mcp.types import ToolAnnotations
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor

def register_connector_tools(app, presentations, get_current_presentation_id, validate_parameters, 
                          is_positive, is_non_negative, is_in_range, is_valid_rgb):
    """Register connector tools with the FastMCP app."""
    
    @app.tool(
        annotations=ToolAnnotations(
            title="Add Connector",
        ),
    )
    def add_connector(
        slide_index: int,
        connector_type: str,
        start_x: float,
        start_y: float,
        end_x: float,
        end_y: float,
        line_width: float = 1.0,
        color: List[int] = None,
        presentation_id: str = None
    ) -> Dict:
        """
        Add connector lines/arrows between points on a slide.
        
        Args:
            slide_index: Index of the slide (0-based)
            connector_type: Type of connector ("straight", "elbow", "curved")
            start_x: Starting X coordinate in inches
            start_y: Starting Y coordinate in inches
            end_x: Ending X coordinate in inches  
            end_y: Ending Y coordinate in inches
            line_width: Width of the connector line in points
            color: RGB color as [r, g, b] list
            presentation_id: Optional presentation ID (uses current if not provided)
            
        Returns:
            Dictionary with operation results
        """
        try:
            # Get presentation
            pres_id = presentation_id or get_current_presentation_id()
            if pres_id not in presentations:
                return {"error": "Presentation not found"}
            
            pres = presentations[pres_id]
            
            # Validate slide index
            if not (0 <= slide_index < len(pres.slides)):
                return {"error": f"Slide index {slide_index} out of range"}
            
            slide = pres.slides[slide_index]
            
            # Map connector types
            connector_map = {
                'straight': MSO_CONNECTOR.STRAIGHT,
                'elbow': MSO_CONNECTOR.ELBOW,
                'curved': MSO_CONNECTOR.CURVED
            }
            
            if connector_type.lower() not in connector_map:
                return {"error": f"Invalid connector type. Use: {list(connector_map.keys())}"}
            
            # Add connector
            connector = slide.shapes.add_connector(
                connector_map[connector_type.lower()],
                Inches(start_x), Inches(start_y),
                Inches(end_x), Inches(end_y)
            )
            
            # Apply formatting
            if line_width:
                connector.line.width = Pt(line_width)
            
            if color and is_valid_rgb(color):
                connector.line.color.rgb = RGBColor(*color)
            
            return {
                "message": f"Added {connector_type} connector to slide {slide_index}",
                "connector_type": connector_type,
                "start_point": [start_x, start_y],
                "end_point": [end_x, end_y],
                "shape_index": len(slide.shapes) - 1
            }
            
        except Exception as e:
            return {"error": f"Failed to add connector: {str(e)}"}