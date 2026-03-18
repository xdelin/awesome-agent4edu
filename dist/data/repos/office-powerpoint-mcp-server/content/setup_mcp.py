# Import necessary Python standard libraries
import os          # For operating with file system, handling files and directory paths
import json        # For processing JSON format data
import subprocess  # For creating and managing subprocesses
import sys         # For accessing Python interpreter related variables and functions
import platform    # For getting current operating system information
import shutil      # For checking if executables exist in PATH

def check_prerequisites():
    """
    Check if necessary prerequisites are installed
    
    Returns:
        tuple: (python_ok, uv_installed, uvx_installed, ppt_server_installed)
    """
    # Check Python version
    python_version = sys.version_info
    python_ok = python_version.major >= 3 and python_version.minor >= 6
    
    # Check if uv/uvx is installed
    uv_installed = shutil.which("uv") is not None
    uvx_installed = shutil.which("uvx") is not None
    
    # Check if office-powerpoint-mcp-server is already installed via pip
    try:
        result = subprocess.run(
            [sys.executable, "-m", "pip", "show", "office-powerpoint-mcp-server"],
            capture_output=True,
            text=True,
            check=False
        )
        ppt_server_installed = result.returncode == 0
    except Exception:
        ppt_server_installed = False
        
    return (python_ok, uv_installed, uvx_installed, ppt_server_installed)

def setup_venv():
    """
    Function to set up Python virtual environment
    
    Features:
    - Checks if Python version meets requirements (3.6+)
    - Creates Python virtual environment (if it doesn't exist)
    - Installs required dependencies in the newly created virtual environment
    
    No parameters required
    
    Returns: Path to Python interpreter in the virtual environment
    """
    # Check Python version
    python_version = sys.version_info
    if python_version.major < 3 or (python_version.major == 3 and python_version.minor < 6):
        print("Error: Python 3.6 or higher is required.")
        sys.exit(1)
    
    # Get absolute path of the directory containing the current script
    base_path = os.path.abspath(os.path.dirname(__file__))
    # Set virtual environment directory path
    venv_path = os.path.join(base_path, '.venv')
    
    # Determine pip and python executable paths based on operating system
    is_windows = platform.system() == "Windows"
    if is_windows:
        pip_path = os.path.join(venv_path, 'Scripts', 'pip.exe')
        python_path = os.path.join(venv_path, 'Scripts', 'python.exe')
    else:
        pip_path = os.path.join(venv_path, 'bin', 'pip')
        python_path = os.path.join(venv_path, 'bin', 'python')
    
    # Check if virtual environment already exists and is valid
    venv_exists = os.path.exists(venv_path)
    pip_exists = os.path.exists(pip_path)
    
    if not venv_exists or not pip_exists:
        print("Creating new virtual environment...")
        # Remove existing venv if it's invalid
        if venv_exists and not pip_exists:
            print("Existing virtual environment is incomplete, recreating it...")
            try:
                shutil.rmtree(venv_path)
            except Exception as e:
                print(f"Warning: Could not remove existing virtual environment: {e}")
                print("Please delete the .venv directory manually and try again.")
                sys.exit(1)
        
        # Create virtual environment
        try:
            subprocess.run([sys.executable, '-m', 'venv', venv_path], check=True)
            print("Virtual environment created successfully!")
        except subprocess.CalledProcessError as e:
            print(f"Error creating virtual environment: {e}")
            sys.exit(1)
    else:
        print("Valid virtual environment already exists.")
    
    # Double-check that pip exists after creating venv
    if not os.path.exists(pip_path):
        print(f"Error: pip executable not found at {pip_path}")
        print("Try creating the virtual environment manually with: python -m venv .venv")
        sys.exit(1)
    
    # Install or update dependencies
    print("\nInstalling requirements...")
    try:
        # Install mcp package
        subprocess.run([pip_path, 'install', 'mcp[cli]'], check=True)
        # Install python-pptx package
        subprocess.run([pip_path, 'install', 'python-pptx'], check=True)
        
        # Also install dependencies from requirements.txt if it exists
        requirements_path = os.path.join(base_path, 'requirements.txt')
        if os.path.exists(requirements_path):
            subprocess.run([pip_path, 'install', '-r', requirements_path], check=True)
        
        
        print("Requirements installed successfully!")
    except subprocess.CalledProcessError as e:
        print(f"Error installing requirements: {e}")
        sys.exit(1)
    except FileNotFoundError:
        print(f"Error: Could not execute {pip_path}")
        print("Try activating the virtual environment manually and installing requirements:")
        if is_windows:
            print(f".venv\\Scripts\\activate")
        else:
            print("source .venv/bin/activate")
        print("pip install mcp[cli] python-pptx")
        sys.exit(1)
    
    return python_path

def generate_mcp_config_local(python_path):
    """
    Generate MCP configuration for locally installed office-powerpoint-mcp-server
    
    Parameters:
    - python_path: Path to Python interpreter in the virtual environment
    
    Returns: Path to the generated config file
    """
    # Get absolute path of the directory containing the current script
    base_path = os.path.abspath(os.path.dirname(__file__))
    
    # Path to PowerPoint Server script
    server_script_path = os.path.join(base_path, 'ppt_mcp_server.py')
    
    # Path to templates directory
    templates_path = os.path.join(base_path, 'templates')
    
    # Create MCP configuration dictionary
    config = {
        "mcpServers": {
            "ppt": {
                "command": python_path,
                "args": [server_script_path],
                "env": {
                    "PYTHONPATH": base_path,
                    "PPT_TEMPLATE_PATH": templates_path
                }
            }
        }
    }
    
    # Save configuration to JSON file
    config_path = os.path.join(base_path, 'mcp-config.json')
    with open(config_path, 'w') as f:
        json.dump(config, f, indent=2)  # indent=2 gives the JSON file good formatting
    
    return config_path

def generate_mcp_config_uvx():
    """
    Generate MCP configuration for PyPI-installed office-powerpoint-mcp-server using UVX
    
    Returns: Path to the generated config file
    """
    # Get absolute path of the directory containing the current script
    base_path = os.path.abspath(os.path.dirname(__file__))
    
    # Path to templates directory (optional for UVX installs)
    templates_path = os.path.join(base_path, 'templates')
    
    # Create MCP configuration dictionary
    env_config = {}
    if os.path.exists(templates_path):
        env_config["PPT_TEMPLATE_PATH"] = templates_path
    
    config = {
        "mcpServers": {
            "ppt": {
                "command": "uvx",
                "args": ["--from", "office-powerpoint-mcp-server", "ppt_mcp_server"],
                "env": env_config
            }
        }
    }
    
    # Save configuration to JSON file
    config_path = os.path.join(base_path, 'mcp-config.json')
    with open(config_path, 'w') as f:
        json.dump(config, f, indent=2)  # indent=2 gives the JSON file good formatting
    
    return config_path

def generate_mcp_config_module():
    """
    Generate MCP configuration for PyPI-installed office-powerpoint-mcp-server using Python module
    
    Returns: Path to the generated config file
    """
    # Get absolute path of the directory containing the current script
    base_path = os.path.abspath(os.path.dirname(__file__))
    
    # Path to templates directory (optional for module installs)
    templates_path = os.path.join(base_path, 'templates')
    
    # Create MCP configuration dictionary
    env_config = {}
    if os.path.exists(templates_path):
        env_config["PPT_TEMPLATE_PATH"] = templates_path
    
    config = {
        "mcpServers": {
            "ppt": {
                "command": sys.executable,
                "args": ["-m", "office_powerpoint_mcp_server"],
                "env": env_config
            }
        }
    }
    
    # Save configuration to JSON file
    config_path = os.path.join(base_path, 'mcp-config.json')
    with open(config_path, 'w') as f:
        json.dump(config, f, indent=2)  # indent=2 gives the JSON file good formatting
    
    return config_path

def install_from_pypi():
    """
    Install office-powerpoint-mcp-server from PyPI
    
    Returns: True if successful, False otherwise
    """
    print("\nInstalling office-powerpoint-mcp-server from PyPI...")
    try:
        subprocess.run([sys.executable, "-m", "pip", "install", "office-powerpoint-mcp-server"], check=True)
        print("office-powerpoint-mcp-server successfully installed from PyPI!")
        return True
    except subprocess.CalledProcessError:
        print("Failed to install office-powerpoint-mcp-server from PyPI.")
        return False

def print_config_instructions(config_path):
    """
    Print instructions for using the generated config
    
    Parameters:
    - config_path: Path to the generated config file
    """
    print(f"\nMCP configuration has been written to: {config_path}")
    
    with open(config_path, 'r') as f:
        config = json.load(f)
    
    print("\nMCP configuration for Claude Desktop:")
    print(json.dumps(config, indent=2))
    
    # Provide instructions for adding configuration to Claude Desktop configuration file
    if platform.system() == "Windows":
        claude_config_path = os.path.expandvars("%APPDATA%\\Claude\\claude_desktop_config.json")
    else:  # macOS
        claude_config_path = os.path.expanduser("~/Library/Application Support/Claude/claude_desktop_config.json")
    
    print(f"\nTo use with Claude Desktop, merge this configuration into: {claude_config_path}")

def create_package_structure():
    """
    Create necessary package structure and directories
    """
    # Get absolute path of the directory containing the current script
    base_path = os.path.abspath(os.path.dirname(__file__))
    
    # Create __init__.py file
    init_path = os.path.join(base_path, '__init__.py')
    if not os.path.exists(init_path):
        with open(init_path, 'w') as f:
            f.write('# PowerPoint MCP Server')
        print(f"Created __init__.py at: {init_path}")
    
    # Create requirements.txt file
    requirements_path = os.path.join(base_path, 'requirements.txt')
    if not os.path.exists(requirements_path):
        with open(requirements_path, 'w') as f:
            f.write('mcp[cli]\npython-pptx\n')
        print(f"Created requirements.txt at: {requirements_path}")
    
    # Create templates directory for PowerPoint templates
    templates_dir = os.path.join(base_path, 'templates')
    if not os.path.exists(templates_dir):
        os.makedirs(templates_dir)
        print(f"Created templates directory at: {templates_dir}")
        
        # Create a README file in templates directory
        readme_path = os.path.join(templates_dir, 'README.md')
        with open(readme_path, 'w') as f:
            f.write("""# PowerPoint Templates

This directory is for storing PowerPoint template files (.pptx or .potx) that can be used with the MCP server.

## Usage

1. Place your template files in this directory
2. Use the `create_presentation_from_template` tool with the template filename
3. The server will automatically search for templates in this directory

## Supported Formats

- `.pptx` - PowerPoint presentation files
- `.potx` - PowerPoint template files

## Example

```python
# Create presentation from template
result = create_presentation_from_template("company_template.pptx")
```

The server will search for templates in:
- Current directory
- ./templates/ (this directory)
- ./assets/
- ./resources/
""")
        print(f"Created templates README at: {readme_path}")
        
        # Offer to create a sample template
        create_sample = input("\nWould you like to create a sample template for testing? (y/n): ").lower().strip()
        if create_sample in ['y', 'yes']:
            create_sample_template(templates_dir)

def create_sample_template(templates_dir):
    """
    Create a sample PowerPoint template for testing
    
    Parameters:
    - templates_dir: Directory where templates are stored
    """
    try:
        # Import required modules for creating a sample template
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        print("Creating sample template...")
        
        # Create a new presentation
        prs = Presentation()
        
        # Get the title slide layout
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Set title and subtitle
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Sample Company Template"
        subtitle.text = "Professional Presentation Template\nCreated by PowerPoint MCP Server"
        
        # Format title
        title_paragraph = title.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(44)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(31, 73, 125)  # Dark blue
        
        # Format subtitle
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(68, 84, 106)  # Gray blue
            paragraph.alignment = PP_ALIGN.CENTER
        
        # Add a content slide
        content_slide_layout = prs.slide_layouts[1]
        content_slide = prs.slides.add_slide(content_slide_layout)
        
        content_title = content_slide.shapes.title
        content_title.text = "Sample Content Slide"
        
        # Add bullet points to content
        content_placeholder = content_slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.text = "Key Features"
        
        # Add bullet points
        bullet_points = [
            "Professional theme and colors",
            "Custom layouts and placeholders", 
            "Ready for content creation",
            "Compatible with MCP server tools"
        ]
        
        for point in bullet_points:
            p = text_frame.add_paragraph()
            p.text = point
            p.level = 1
        
        # Add a section header slide
        section_slide_layout = prs.slide_layouts[2] if len(prs.slide_layouts) > 2 else prs.slide_layouts[0]
        section_slide = prs.slides.add_slide(section_slide_layout)
        
        if section_slide.shapes.title:
            section_slide.shapes.title.text = "Template Features"
        
        # Save the sample template
        template_path = os.path.join(templates_dir, 'sample_template.pptx')
        prs.save(template_path)
        
        print(f"✅ Sample template created: {template_path}")
        print("   You can now test the template feature with:")
        print("   • get_template_info('sample_template.pptx')")
        print("   • create_presentation_from_template('sample_template.pptx')")
        
    except ImportError:
        print("⚠️  Cannot create sample template: python-pptx not installed yet")
        print("   Run the setup first, then manually create templates in the templates/ directory")
    except Exception as e:
        print(f"❌ Failed to create sample template: {str(e)}")
        print("   You can manually add template files to the templates/ directory")

# Main execution entry point
if __name__ == '__main__':
    # Check prerequisites
    python_ok, uv_installed, uvx_installed, ppt_server_installed = check_prerequisites()
    
    if not python_ok:
        print("Error: Python 3.6 or higher is required.")
        sys.exit(1)
    
    print("PowerPoint MCP Server Setup")
    print("===========================\n")
    
    # Create necessary files
    create_package_structure()
    
    # If office-powerpoint-mcp-server is already installed, offer config options
    if ppt_server_installed:
        print("office-powerpoint-mcp-server is already installed via pip.")
        
        if uvx_installed:
            print("\nOptions:")
            print("1. Generate MCP config for UVX (recommended)")
            print("2. Generate MCP config for Python module")
            print("3. Set up local development environment")
            
            choice = input("\nEnter your choice (1-3): ")
            
            if choice == "1":
                config_path = generate_mcp_config_uvx()
                print_config_instructions(config_path)
            elif choice == "2":
                config_path = generate_mcp_config_module()
                print_config_instructions(config_path)
            elif choice == "3":
                python_path = setup_venv()
                config_path = generate_mcp_config_local(python_path)
                print_config_instructions(config_path)
            else:
                print("Invalid choice. Exiting.")
                sys.exit(1)
        else:
            print("\nOptions:")
            print("1. Generate MCP config for Python module")
            print("2. Set up local development environment")
            
            choice = input("\nEnter your choice (1-2): ")
            
            if choice == "1":
                config_path = generate_mcp_config_module()
                print_config_instructions(config_path)
            elif choice == "2":
                python_path = setup_venv()
                config_path = generate_mcp_config_local(python_path)
                print_config_instructions(config_path)
            else:
                print("Invalid choice. Exiting.")
                sys.exit(1)
    
    # If office-powerpoint-mcp-server is not installed, offer installation options
    else:
        print("office-powerpoint-mcp-server is not installed.")
        
        print("\nOptions:")
        print("1. Install from PyPI (recommended)")
        print("2. Set up local development environment")
        
        choice = input("\nEnter your choice (1-2): ")
        
        if choice == "1":
            if install_from_pypi():
                if uvx_installed:
                    print("\nNow generating MCP config for UVX...")
                    config_path = generate_mcp_config_uvx()
                else:
                    print("\nUVX not found. Generating MCP config for Python module...")
                    config_path = generate_mcp_config_module()
                print_config_instructions(config_path)
        elif choice == "2":
            python_path = setup_venv()
            config_path = generate_mcp_config_local(python_path)
            print_config_instructions(config_path)
        else:
            print("Invalid choice. Exiting.")
            sys.exit(1)
    
    print("\nSetup complete! You can now use the PowerPoint MCP server with compatible clients like Claude Desktop.")
    
    print("\n" + "="*60)
    print("POWERPOINT MCP SERVER - NEW FEATURES")
    print("="*60)
    print("\n📁 Template Support:")
    print("   • Place PowerPoint templates (.pptx/.potx) in the ./templates/ directory")
    print("   • Use 'create_presentation_from_template' tool to create presentations from templates")
    print("   • Use 'get_template_info' tool to inspect template layouts and properties")
    print("   • Templates preserve branding, themes, and custom layouts")
    print("   • Template path configured via PPT_TEMPLATE_PATH environment variable")
    
    print("\n🔧 Available MCP Tools:")
    print("   Presentations:")
    print("   • create_presentation - Create new blank presentation")
    print("   • create_presentation_from_template - Create from template file")
    print("   • get_template_info - Inspect template file details")
    print("   • open_presentation - Open existing presentation")
    print("   • save_presentation - Save presentation to file")
    
    print("\n   Content:")
    print("   • add_slide - Add slides with various layouts")
    print("   • add_textbox - Add formatted text boxes")
    print("   • add_image - Add images from files or base64")
    print("   • add_table - Add formatted tables")
    print("   • add_shape - Add various auto shapes")
    print("   • add_chart - Add column, bar, line, and pie charts")
    
    print("\n📚 Documentation:")
    print("   • Full API documentation available in README.md")
    print("   • Template usage examples included")
    print("   • Check ./templates/README.md for template guidelines")
    
    print("\n🚀 Quick Start with Templates:")
    print("   1. Copy your .pptx template to ./templates/")
    print("   2. Use: create_presentation_from_template('your_template.pptx')")
    print("   3. Add slides using template layouts")
    print("   4. Save your presentation")
    print("\n💡 Custom Template Paths:")
    print("   • Set PPT_TEMPLATE_PATH environment variable for custom locations")
    print("   • Supports multiple paths (colon-separated on Unix, semicolon on Windows)")
    print("   • Example: PPT_TEMPLATE_PATH='/path/to/templates:/path/to/more/templates'")
    
    print("\n" + "="*60)