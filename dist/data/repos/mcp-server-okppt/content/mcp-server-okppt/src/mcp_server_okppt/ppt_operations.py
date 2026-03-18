from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from typing import Optional, Union, List, Dict, Any # Added Dict, Any for type hinting
import os
import datetime
import traceback
import logging # Added for logging
import sys
# Assuming logger is configured in server.py and can be imported or a new one is created here
# For simplicity, let's assume a logger is available or we'll create one.
logger = logging.getLogger(__name__)
if not logger.handlers:
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
        stream=sys.stdout
    )

# Path helper functions (will be imported from server.py or defined if not available)
# For now, we'll assume they are available from server.py
# from server import get_default_output_path, cleanup_filename, get_output_dir, get_tmp_dir

def get_placeholder_type_name(ph_type):
    """将占位符类型ID转换为可读名称"""
    placeholder_types = {
        1: "标题 (Title)",
        2: "正文/内容 (Body)", 
        3: "居中标题 (Center Title)",
        4: "副标题 (Subtitle)",
        5: "日期和时间 (Date)",
        6: "幻灯片编号 (Slide Number)",
        7: "页脚 (Footer)",
        8: "页眉 (Header)",
        9: "对象 (Object)",
        10: "图表 (Chart)",
        11: "表格 (Table)",
        12: "剪贴画 (Clip Art)",
        13: "组织结构图 (Organization Chart)",
        14: "媒体剪辑 (Media Clip)",
        15: "图片 (Picture)",
        16: "垂直对象 (Vertical Object)",
        17: "垂直文本 (Vertical Text)"
    }
    return placeholder_types.get(ph_type, f"未知类型 (Unknown Type {ph_type})")

# ... rest of the functions will be added below ... 

# Path helper functions will be imported from server.py
# For now, we assume they are available.
# It's better to pass them or have server.py provide them.
# For this step, direct import for get_default_output_path for simplicity if needed, though analyze_layout_details doesn't create files.
try:
    from mcp_server_okppt.server import get_default_output_path, cleanup_filename
except ImportError:
    logger.warning("Could not import path helpers from server.py for ppt_operations.py")
    def get_default_output_path(file_type="pptx", base_name=None, op_type=None): # Dummy for standalone
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        return f"{base_name or 'temp'}_{op_type or 'op'}_{timestamp}.{file_type}"
    def cleanup_filename(filename: str) -> str:
        return filename


def analyze_layout_details(prs_path: str, title: str = "演示文稿") -> Dict[str, Any]:
    """分析演示文稿的布局详细信息并返回结构化数据"""
    analysis_data: Dict[str, Any] = {
        "presentation_path": prs_path,
        "analysis_title": title,
        "slide_count": 0,
        "master_count": 0,
        "total_layouts_count": 0,
        "layout_type_stats": {},
        "masters_details": [],
        "slide_layout_usage_summary": {},
        "slides_details": [],
        "unused_layouts_summary": [],
        "layout_utilization": {}
    }

    try:
        prs = Presentation(prs_path)
        analysis_data["slide_count"] = len(prs.slides)
        analysis_data["master_count"] = len(prs.slide_masters)

        total_layouts = 0
        layout_type_stats: Dict[str, int] = {
            "自定义布局": 0,
            "系统布局": 0,
            "无名称布局": 0
        }
        
        masters_details_list = []
        for master_idx, master in enumerate(prs.slide_masters):
            master_detail = {
                "master_index": master_idx + 1,
                "master_name": master.name,
                "layout_count": len(master.slide_layouts),
                "layouts": []
            }
            total_layouts += len(master.slide_layouts)
            
            for layout_idx, layout in enumerate(master.slide_layouts):
                layout_name = layout.name
                layout_type = ""
                layout_display = ""

                if layout_name == "" or layout_name is None:
                    layout_type = "无名称布局"
                    layout_display = f"无名称（默认布局）"
                elif layout_name.isdigit():
                    layout_type = "系统布局"
                    layout_display = f"系统布局 - {layout_name}"
                else:
                    layout_type = "自定义布局"
                    layout_display = f"自定义布局 - {layout_name}"
                
                layout_type_stats[layout_type] = layout_type_stats.get(layout_type, 0) + 1
                
                layout_detail = {
                    "layout_index": layout_idx + 1,
                    "layout_name_original": layout_name,
                    "layout_display_name": layout_display,
                    "layout_type": layout_type,
                    "placeholder_count": len(layout.placeholders),
                    "placeholders": []
                }
                
                if len(layout.placeholders) > 0:
                    for ph_idx, placeholder in enumerate(layout.placeholders):
                        try:
                            ph_format = placeholder.placeholder_format
                            ph_type_code = ph_format.type
                            ph_access_id = ph_format.idx
                            ph_type_name = get_placeholder_type_name(ph_type_code)
                            layout_detail["placeholders"].append({
                                "placeholder_index": ph_idx + 1,
                                "type_name": ph_type_name,
                                "access_id": ph_access_id,
                                "type_code": ph_type_code
                            })
                        except Exception as e_ph:
                            logger.warning(f"Placeholder {ph_idx + 1} in layout '{layout_name}': type无法获取 ({str(e_ph)})")
                            layout_detail["placeholders"].append({
                                "placeholder_index": ph_idx + 1,
                                "error": f"类型无法获取 ({str(e_ph)})"
                            })
                master_detail["layouts"].append(layout_detail)
            masters_details_list.append(master_detail)
        
        analysis_data["masters_details"] = masters_details_list
        analysis_data["total_layouts_count"] = total_layouts
        analysis_data["layout_type_stats"] = {
            k: {"count": v, "percentage": (v / total_layouts * 100) if total_layouts > 0 else 0}
            for k, v in layout_type_stats.items()
        }

        slides_details_list = []
        layout_usage: Dict[str, int] = {}
        if len(prs.slides) > 0:
            for slide_idx, slide in enumerate(prs.slides):
                layout_name = slide.slide_layout.name
                layout_type = ""
                layout_display = ""

                if layout_name == "" or layout_name is None:
                    layout_type = "无名称布局"
                    layout_display = "无名称（默认布局）"
                elif layout_name.isdigit():
                    layout_type = "系统布局"
                    layout_display = f"系统布局 - {layout_name}"
                else:
                    layout_type = "自定义布局"
                    layout_display = f"自定义布局 - {layout_name}"
                
                slide_title = "无标题"
                try:
                    if slide.shapes.title and slide.shapes.title.has_text_frame and slide.shapes.title.text.strip():
                        slide_title = slide.shapes.title.text.strip()
                except: # pylint: disable=bare-except
                    pass 
                
                slide_detail = {
                    "slide_number": slide_idx + 1,
                    "title": slide_title,
                    "used_layout_name_original": layout_name,
                    "used_layout_display_name": layout_display,
                    "used_layout_type": layout_type,
                    "placeholder_count_on_slide": len(slide.placeholders)
                }
                slides_details_list.append(slide_detail)
                
                layout_usage[layout_name] = layout_usage.get(layout_name, 0) + 1
            
            analysis_data["slide_layout_usage_summary"] = {
                name: {"count": count, "percentage": (count / len(prs.slides) * 100)}
                for name, count in sorted(layout_usage.items(), key=lambda x: x[1], reverse=True)
            }

            unused_layouts_list = []
            for master in prs.slide_masters:
                for layout in master.slide_layouts:
                    if layout.name not in layout_usage:
                        layout_type_unused = "自定义布局"
                        if layout.name == "" or layout.name is None: layout_type_unused = "无名称布局"
                        elif layout.name.isdigit(): layout_type_unused = "系统布局"
                        unused_layouts_list.append({"name": layout.name, "type": layout_type_unused})
            analysis_data["unused_layouts_summary"] = unused_layouts_list
            
            used_layouts_count = len(layout_usage)
            total_available_layouts = sum(len(master.slide_layouts) for master in prs.slide_masters)
            utilization_rate = (used_layouts_count / total_available_layouts * 100) if total_available_layouts > 0 else 0
            analysis_data["layout_utilization"] = {
                "total_available": total_available_layouts,
                "used_count": used_layouts_count,
                "utilization_rate_percentage": utilization_rate
            }
        analysis_data["slides_details"] = slides_details_list

        return {
            "status": "success",
            "message": "Presentation analysis successful.",
            "data": analysis_data,
            "output_path": None
        }

    except Exception as e:
        logger.error(f"Error in analyze_layout_details for {prs_path}: {str(e)}\n{traceback.format_exc()}")
        return {
            "status": "error",
            "message": f"分析布局详情失败: {str(e)}",
            "data": analysis_data, # Return partial data if any
            "output_path": None
        }

def insert_layout(prs_path: str, layout_to_insert: str, output_path: Optional[str] = None, new_slide_title: Optional[str] = None) -> Dict[str, Any]:
    """
    在演示文稿中插入一个使用指定布局的新幻灯片，并返回操作结果。
    """
    prs = None
    original_slide_count = 0
    try:
        prs = Presentation(prs_path)
        original_slide_count = len(prs.slides)

        base_name_cleaned = cleanup_filename(os.path.splitext(os.path.basename(prs_path))[0])
        final_output_path = output_path or get_default_output_path(
            base_name=base_name_cleaned,
            op_type=f"inserted_layout_{layout_to_insert.replace(' ', '_')}"
        )
        
        available_layouts = {}
        for master in prs.slide_masters:
            for layout in master.slide_layouts:
                available_layouts[layout.name] = layout
        
        target_layout = available_layouts.get(layout_to_insert)
        
        if target_layout:
            new_slide = prs.slides.add_slide(target_layout)
            logger.info(f"Added new slide using layout '{layout_to_insert}'")
            
            if new_slide_title:
                if new_slide.shapes.title:
                    new_slide.shapes.title.text = new_slide_title
                    logger.info(f"New slide title set to: '{new_slide_title}'")
                else:
                    logger.warning(f"Layout '{layout_to_insert}' has no title placeholder for title '{new_slide_title}'")
            
            prs.save(final_output_path)
            logger.info(f"Presentation saved to {final_output_path}")
            
            return {
                "status": "success",
                "message": f"Successfully inserted slide with layout '{layout_to_insert}'.",
                "data": {
                    "slides_total": len(prs.slides),
                    "original_slide_count": original_slide_count,
                    "new_slide_title_set": new_slide_title if new_slide.shapes.title and new_slide_title else None
                },
                "output_path": final_output_path
            }
        else:
            logger.error(f"Layout '{layout_to_insert}' not found in {prs_path}. Available: {list(available_layouts.keys())}")
            return {
                "status": "error",
                "message": f"Layout '{layout_to_insert}' not found. Available layouts: {', '.join(available_layouts.keys())}",
                "data": {"available_layouts": list(available_layouts.keys()), "original_slide_count": original_slide_count},
                "output_path": None
            }
            
    except Exception as e:
        logger.error(f"Error in insert_layout for {prs_path} with layout '{layout_to_insert}': {str(e)}\n{traceback.format_exc()}")
        return {
            "status": "error",
            "message": f"插入布局失败: {str(e)}",
            "data": {"original_slide_count": original_slide_count if prs else 0},
            "output_path": None
        }

def clear_placeholder_content(prs_path: str, output_path: Optional[str] = None, slide_indices: Optional[List[int]] = None) -> Dict[str, Any]:
    """
    清空演示文稿中占位符的内容（保留占位符结构），并返回操作结果。
    """
    try:
        prs = Presentation(prs_path)
        original_slide_count = len(prs.slides)

        base_name_cleaned = cleanup_filename(os.path.splitext(os.path.basename(prs_path))[0])
        final_output_path = output_path or get_default_output_path(
            base_name=base_name_cleaned,
            op_type="content_cleared"
        )
        
        total_cleared = 0
        slides_processed = 0
        processed_slides_info = []

        slides_to_process_indices = list(range(len(prs.slides))) if slide_indices is None else [i for i in slide_indices if 0 <= i < len(prs.slides)]
        
        logger.info(f"Preparing to clear content from {len(slides_to_process_indices)} slides in {prs_path}.")

        for slide_idx in slides_to_process_indices:
            slide = prs.slides[slide_idx]
            slide_cleared_count = 0
            placeholders_on_slide_info = []
            
            for ph_idx, placeholder in enumerate(slide.placeholders):
                try:
                    placeholder_cleared_flag = False
                    ph_access_id = placeholder.placeholder_format.idx
                    ph_type_name = get_placeholder_type_name(placeholder.placeholder_format.type)
                    
                    current_ph_info = {"access_id": ph_access_id, "type": ph_type_name, "cleared": False, "reason": ""}

                    if hasattr(placeholder, 'text_frame') and placeholder.text_frame is not None:
                        if placeholder.text_frame.text.strip():
                            placeholder.text_frame.clear()
                            placeholder_cleared_flag = True
                            current_ph_info["reason"] = "text_frame cleared"
                            logger.debug(f"  Cleared text_frame for placeholder {ph_access_id} on slide {slide_idx + 1}")
                    elif hasattr(placeholder, 'text'): # pylint: disable=else-if-used
                        if placeholder.text.strip(): 
                            placeholder.text = ""
                            placeholder_cleared_flag = True
                            current_ph_info["reason"] = "text property cleared"
                            logger.debug(f"  Cleared text property for placeholder {ph_access_id} on slide {slide_idx + 1}")
                    
                    if placeholder_cleared_flag:
                        slide_cleared_count += 1
                        total_cleared += 1
                        current_ph_info["cleared"] = True
                    else:
                        current_ph_info["reason"] = "no text content or not clearable type"

                    placeholders_on_slide_info.append(current_ph_info)
                
                except Exception as e_ph:
                    logger.warning(f"  Error clearing placeholder {ph_idx} (ID: {ph_access_id if 'ph_access_id' in locals() else 'N/A'}) on slide {slide_idx + 1}: {str(e_ph)}")
                    placeholders_on_slide_info.append({"access_id": ph_access_id if 'ph_access_id' in locals() else 'N/A', "type": ph_type_name if 'ph_type_name' in locals() else 'N/A', "cleared": False, "reason": f"error: {str(e_ph)}"})
            
            processed_slides_info.append({
                "slide_number": slide_idx + 1,
                "cleared_count_on_slide": slide_cleared_count,
                "placeholders_status": placeholders_on_slide_info
            })
            if slide_cleared_count > 0:
                slides_processed += 1
        
        prs.save(final_output_path)
        logger.info(f"Presentation with cleared content saved to {final_output_path}")
        
        return {
            "status": "success",
            "message": f"Successfully cleared content from {total_cleared} placeholder(s) in {slides_processed} slide(s).",
            "data": {
                "slides_processed_count": slides_processed,
                "placeholders_cleared_total": total_cleared,
                "slides_targetted_count": len(slides_to_process_indices),
                "processed_slides_details": processed_slides_info
            },
            "output_path": final_output_path
        }
        
    except Exception as e:
        logger.error(f"Error in clear_placeholder_content for {prs_path}: {str(e)}\n{traceback.format_exc()}")
        return {
            "status": "error",
            "message": f"清空占位符内容失败: {str(e)}",
            "data": {"placeholders_cleared_total": total_cleared if 'total_cleared' in locals() else 0, "slides_processed_count": slides_processed if 'slides_processed' in locals() else 0},
            "output_path": None
        }

def assign_placeholder_content(prs_path: str, slide_index: int, placeholder_access_id: int, content: str, output_path: Optional[str] = None) -> Dict[str, Any]:
    """
    给演示文稿中特定幻灯片的特定占位符赋值，并返回操作结果。
    """
    try:
        prs = Presentation(prs_path)

        if not (0 <= slide_index < len(prs.slides)):
            message = f"幻灯片索引 {slide_index} 超出范围 (共 {len(prs.slides)} 张幻灯片)。"
            logger.error(message)
            return {"status": "error", "message": message, "data": None, "output_path": None}
        
        slide = prs.slides[slide_index]
        target_placeholder = None
        
        try:
            # Attempt to access placeholder by integer index first if access_id is int
            if isinstance(placeholder_access_id, int):
                 # Check if the direct index access is valid for the slide's placeholders list
                if 0 <= placeholder_access_id < len(slide.placeholders):
                    # This assumes placeholder_access_id might be used as a direct 0-based index
                    # However, python-pptx usually relies on placeholder_format.idx for uniqueness if available
                    # For robustness, we should iterate or use a more reliable way if idx is not a direct list index.
                    pass # This path might need more thought if access_id isn't a direct list index

            # Standard way: iterate and check placeholder_format.idx
            found_by_idx_check = False
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == placeholder_access_id:
                    target_placeholder = ph
                    found_by_idx_check = True
                    break
            
            if not found_by_idx_check:
                # Fallback: if placeholder_access_id was meant as a direct 0-based index for slide.placeholders list
                # This is less robust as placeholder_format.idx is the true unique ID from the layout.
                if isinstance(placeholder_access_id, int) and 0 <= placeholder_access_id < len(slide.placeholders):
                    target_placeholder = slide.placeholders[placeholder_access_id]
                    logger.warning(f"Accessed placeholder on slide {slide_index + 1} by direct list index {placeholder_access_id} as fallback. "
                                   f"Prefer using placeholder_format.idx for 'placeholder_access_id'.")
                else:
                    raise KeyError # Trigger the common error handling below

        except KeyError:
            available_phs_info = []
            for ph_item in slide.placeholders:
                available_phs_info.append({
                    "access_id": ph_item.placeholder_format.idx,
                    "type": get_placeholder_type_name(ph_item.placeholder_format.type)
                })
            message = f"在幻灯片 {slide_index + 1} 上未找到访问ID为 {placeholder_access_id} 的占位符。"
            logger.error(f"{message} Layout: '{slide.slide_layout.name}'. Available IDs: {available_phs_info}")
            return {"status": "error", "message": message, "data": {"available_placeholders": available_phs_info}, "output_path": None}
        
        if target_placeholder is None: # Should be caught by KeyError, but as a safeguard
            message = f"无法通过访问ID {placeholder_access_id} 找到占位符 (safeguard check)."
            logger.error(message)
            return {"status": "error", "message": message, "data": None, "output_path": None}

        logger.info(f"Successfully found placeholder with access_id {placeholder_access_id} on slide {slide_index+1}. Type: {get_placeholder_type_name(target_placeholder.placeholder_format.type)}")

        assigned = False
        assignment_method = "none"
        if hasattr(target_placeholder, 'text_frame') and target_placeholder.text_frame is not None:
            target_placeholder.text_frame.text = str(content)
            assigned = True
            assignment_method = "text_frame"
            logger.info(f"  Assigned content to placeholder's text_frame.")
        elif hasattr(target_placeholder, 'text'):
            target_placeholder.text = str(content)
            assigned = True
            assignment_method = "text_property"
            logger.info(f"  Assigned content to placeholder's text property.")
        else:
            ph_type_name = get_placeholder_type_name(target_placeholder.placeholder_format.type)
            logger.warning(f"Placeholder type '{ph_type_name}' (ID: {placeholder_access_id}) might not support direct text assignment.")

        if not assigned:
            message = f"未能成功赋值内容到占位符ID {placeholder_access_id} (类型: {get_placeholder_type_name(target_placeholder.placeholder_format.type)})."
            logger.error(message)
            return {"status": "error", "message": message, "data": {"placeholder_type": get_placeholder_type_name(target_placeholder.placeholder_format.type)}, "output_path": None}

        base_name_cleaned = cleanup_filename(os.path.splitext(os.path.basename(prs_path))[0])
        final_output_path = output_path or get_default_output_path(
            base_name=base_name_cleaned, 
            op_type=f"assigned_S{slide_index}_P{placeholder_access_id}"
        )
        
        prs.save(final_output_path)
        logger.info(f"Presentation with assigned content saved to {final_output_path}")

        return {
            "status": "success",
            "message": f"Successfully assigned content to placeholder ID {placeholder_access_id} on slide {slide_index + 1}.",
            "data": {
                "slide_index": slide_index,
                "placeholder_access_id": placeholder_access_id,
                "content_assigned_length": len(str(content)),
                "assignment_method": assignment_method
            },
            "output_path": final_output_path
        }

    except Exception as e:
        logger.error(f"Error in assign_placeholder_content for {prs_path}, slide {slide_index}, ph_id {placeholder_access_id}: {str(e)}\n{traceback.format_exc()}")
        return {
            "status": "error",
            "message": f"给占位符赋值时发生严重错误: {str(e)}",
            "data": None,
            "output_path": None
        }

# ... other functions will follow (placeholder for future, or end of relevant functions)

# ... other functions will follow ... 