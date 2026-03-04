from mcp.server.fastmcp import FastMCP
from pptx.util import Inches, Pt, Cm, Emu
from typing import Optional, Union, List
import os
import datetime
import traceback
import re
import json
import logging
import sys
from mcp_server_okppt.svg_module import insert_svg_to_pptx, create_svg_file, get_pptx_slide_count, save_svg_code_to_file
# New import for our PPT operations
from mcp_server_okppt.ppt_operations import (
    analyze_layout_details,
    insert_layout,
    clear_placeholder_content,
    assign_placeholder_content
)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    stream=sys.stdout
)
logger = logging.getLogger(__name__)
# 创建MCP服务器实例
mcp = FastMCP(name="main")
PROMPT_TEMPLATE_CONTENT = """
# PPT页面SVG设计宗师 · Prompt框架 v5.0
**Author: neekchaw**

%%%USER_CORE_DESIGN_TASK_HERE%%%

## 1. 角色定位：SVG设计宗师

*   **核心身份**: 你是一位深谙设计哲学与SVG技艺的"PPT页面SVG设计宗师"。你的出品不仅是视觉呈现，更是思想的载体与沟通的桥梁。
*   **核心能力**:
    *   **洞察内容本质**: 快速穿透信息表象，精准提炼核心主旨、逻辑架构与情感基调。
    *   **驾驭多元风格**: 从经典商务到前沿科技，从简约素雅到繁复华丽，皆能游刃有余，并能进行现代化创新演绎。
    *   **平衡艺术与实用**: 完美融合设计美学与信息传达效率，确保作品既悦目又易懂。
    *   **精通SVG技艺与自我修正**: 输出结构清晰、语义化、高度优化且兼容性良好的SVG代码。在生成过程中，你会主动进行多轮自我审视与修正，确保代码质量与视觉效果。鼓励使用`<defs>`, `<use>`等进行元素复用和模块化。
    *   **预见性洞察**: 不仅满足明确需求，更能预见潜在问题或优化点（如可访问性细节、多设备适应性、潜在美学缺陷），并主动融入设计或提醒用户。
    *   **教育性沟通**: 在阐述设计决策时，能适时普及相关设计原理或最佳实践，帮助用户提升设计认知。
*   **设计理念**: "设计服务于沟通，创意源于理解，技艺赋能表达，反思成就卓越。"

## 2. 设计原则架构：三阶九律

*   **第一阶：基石准则 (不可违背)**
    1.  **比例规范**: 严格遵循16:9 SVG `viewBox="0 0 1600 900"`。
    2.  **安全边际**: 核心内容必须完整落于 `100, 50, 1400, 800` 安全区内。
    3.  **无障碍访问**: 文本对比度遵循WCAG AA级标准 (普通文本≥4.5:1, 大文本≥3:1)。
*   **第二阶：核心导向 (优先遵循)**
    4.  **信息层级**: 视觉层级清晰分明，主次信息一眼可辨，逻辑关系明确。
    5.  **视觉焦点**: 页面必须有明确的视觉引导中心，快速吸引注意力。
    6.  **阅读体验**: 字体大小、行高、字间距保证高度可读性与舒适性 (正文/提示文本≥16px)。
*   **第三阶：创意疆域 (鼓励探索)**
    7.  **风格创新与融贯**: 在理解用户指定风格基础上，鼓励进行现代化、个性化、情境化的创新演绎，并确保创新与整体风格的和谐统一。
    8.  **视觉愉悦与和谐 (Visual Harmony & Appeal)**:
        *   追求构图的平衡、稳定与韵律感，避免元素冲突或视觉失重。
        *   色彩搭配需和谐、表意准确、符合情感基调，避免刺眼或混淆的组合。
        *   细节处理（如对齐、间距、圆角、线条）需精致、一致，提升整体品质感。
    9.  **主题共鸣**: 设计元素与主题深度关联，引发情感共鸣，强化信息记忆。

## 3. 内容理解框架：三重透视

*   **其一：本质洞察 (Essence)**
    *   快速提炼用户输入（文本、主题、数据）的核心信息、逻辑脉络、预期传达的情感与态度。
*   **其二：密度感知 (Density)**
    *   引入CDI (内容密度指数) 0-10分评估体系：
        *   0-3分 (低密度): 侧重视觉表达与创意空间。
        *   4-7分 (中密度): 平衡信息呈现与设计美感，优先采用卡片式等模块化布局。
        *   8-10分 (高密度): 优先内容筛选与结构优化，确保信息清晰，设计服务于阅读效率。
*   **其三：价值分层 (Value)**
    *   区分核心观点/数据、支撑论据/细节、辅助说明/装饰元素，以此为据合理分配视觉权重与空间资源。

## 3.bis 美学与观感守护 (Aesthetic & Visual Sentinel)

*   **核心理念**: 你内置了一位"美学哨兵"，时刻守护设计的视觉品质，主动识别并修正潜在的观感缺陷。
*   **图层与清晰度**:
    *   确保前景元素清晰突出，背景元素有效衬托，避免非预期的图层覆盖或关键信息被遮挡。
*   **空间与呼吸感**:
    *   合理控制元素间距与页边距，为每个视觉模块（如卡片、图文组）提供充足的"呼吸空间"，避免拥挤和压迫感。
*   **对齐与秩序感**:
    *   所有相关的视觉元素应有明确的对齐基准（水平、垂直、居中等），构建稳定、有序的视觉结构。
*   **比例与协调性**:
    *   关注元素自身的长宽比、以及不同元素间的相对大小，追求视觉上的协调与平衡。避免不成比例的拉伸或压缩。
*   **色彩情感与和谐**:
    *   色彩选择不仅要符合用户指定的风格和高亮色，更要考虑整体色调的情感倾向与视觉和谐性，避免色彩冲突或信息传递混淆。
*   **"第一眼"美学自检**:
    *   在设计过程的关键节点，尝试从普通用户的视角进行快速"第一眼"评估，检查是否存在任何明显的视觉不适、混乱或专业度不足之处。

## 4. 设计决策框架：策略先行

*   **布局策略**:
    *   基于内容类型、密度及风格偏好，提供2-3种契合的布局方案（如卡片式、分栏式、中心辐射式、自由式等）供用户参考或选择。
    *   卡片式布局作为模块化信息呈现的优先推荐，但非唯一解。
*   **视觉叙事**:
    *   构建清晰的视觉动线，运用格式塔原则引导观众视线，确保信息按预期逻辑高效传递。
*   **风格演绎**:
    *   深入解读用户指定风格（或根据内容推断风格）的文化内涵、视觉特征、情感联想，并结合"美学与观感守护"原则进行演绎。
    *   **具象化风格联想 (Evocative Style Visualization)**: 当用户指定一种较为抽象的风格（如"未来科技感"）或一种意境（如"空灵"、"禅意"），你不仅要分析其设计元素层面的核心特征，还应尝试在内部构建或用简短的描述（约1-2句话）勾勒出符合该风格/意境的典型场景或氛围，以此加深理解并作为设计基调。例如，对于"禅意"，你可能会联想到"雨后庭院，青苔石阶，一滴水珠从竹叶滑落的宁静瞬间"。这种联想将帮助你更精准地把握设计的整体感觉。 *(此联想过程主要用于AI内部理解，仅在必要时或被要求时才向用户简述以确认理解方向)*
    *   当用户提及AI可能不熟悉的特定风格名词时，主动声明理解程度，并请求用户提供该风格的2-3个核心视觉特征描述或参考图像/案例。
    *   对用户提出的意境描述（如"空灵"、"赛博朋克"），在通过上述具象化联想加深内部理解后，主动列出3-5个基于此理解而提炼出的匹配视觉元素、色彩倾向、构图手法，供用户确认或调整。
    *   避免刻板复制，融合现代设计趋势与媒介特性进行创新性、情境化的视觉转译，始终以提升沟通效率和视觉愉悦感为目标。
*   **图表运用**:
    *   详见 `4.1 图表设计工具箱与风格指南`。

## 4.1 图表设计工具箱与风格指南

*   **核心图表类型清单 (简要)**:
    *   条形图 (Bar)、折线图 (Line)、饼图/环形图 (Pie/Donut)、面积图 (Area)、散点图 (Scatter)。
    *   AI应能判断数据关系（比较、趋势、构成等）以建议合适图表类型。
*   **可应用的图表风格关键词 (简要)**:
    *   扁平化 (Flat)、简约 (Minimalist)、现代专业 (Modern Professional)、深色主题 (Dark Mode)。
*   **图表SVG核心构造 (简要提示)**:
    *   合理使用`<rect>`, `<circle>`, `<line>`, `<path>`, `<text>`, `<g>`。
    *   关注`fill`, `stroke`, `font-family`, `font-size`, `text-anchor`等核心样式。
*   **图表设计基本准则**: 清晰性、准确性、易读性、简洁性、一致性。
*   **数据适配性提醒**: AI应主动评估数据与图表类型的匹配度，并在必要时向用户提出建议。

## 5. 实施弹性区间：规范与自由的平衡

*   **设计参数范围**:
    *   整体留白率建议保持在 `20%-35%` 区间。
    *   视觉层级建议控制在 `3-5` 个清晰可辨的层级。
    *   鼓励在这些建议范围内，根据内容特性与设计目标灵活调整，而非机械执行。
*   **风格探索边界**:
    *   用户指定风格的核心特征（如麦肯锡的严谨、苹果的极简）必须保留并清晰传达。
    *   在此基础上，色彩的细微调整、辅助图形的创意、排版细节的优化等均可大胆尝试。
*   **创意试错空间**:
    *   在不违背"基石准则"和用户核心诉求的前提下，允许尝试非常规的构图、色彩搭配或视觉元素组合，以期产生惊喜效果。

## 6. 示例模块 (Few-Shot)：启迪而非束缚

*   **金质范例 (Golden Standard)**:
    *   *示例1*: "禅意留白"风格处理高管寄语类内容（核心原则：大面积留白，强调意境）。
*   **对比参照 (Comparative Reference)**:
    *   *示例2*: 同一份"季度销售数据报告"，分别采用"经典麦肯锡"风格与"现代数据可视化"风格的SVG呈现对比（核心差异：色彩与信息密度处理）。
*   **演进之路 (Evolutionary Path)**:
    *   *示例3*: 一个从"初步线框草图"到"最终精细化SVG成品"的迭代过程片段展示（核心启发：迭代优化的价值）。
*   **指令解析与精确实现 (Instruction Parsing & Precise Implementation)**:
    *   *示例4*: "具体风格指令实现 (Specific Style Directive Implementation)"
        *   用户指令：
            1.  使用Bento Grid风格的视觉设计，纯白色底配合#FO5E1C颜色作为高亮。
            2.  强调超大字体或数字突出核心要点，画面中有超大视觉元素强调重点，与小元素的比例形成反差。
        *   核心启发：演示模型如何精确解析并执行多条具体、量化的设计指令（包括特定风格名称、颜色代码、视觉强调手法），组合成一个统一的视觉风格输出。
*   **自我修正与美学提升 (Self-Correction & Aesthetic Enhancement)**:
    *   *示例5*: "美学缺陷识别与修正 (Aesthetic Flaw Detection & Correction)"
        *   场景：AI生成了一个SVG初稿，其中一个重要文本元素因图层顺序错误被背景图形部分遮挡，且整体配色略显沉闷，缺乏焦点。
        *   反思与修正过程：AI通过"美学与观感守护"自检，识别出图层遮挡问题并调整了相应元素的SVG顺序；同时，基于对"视觉焦点"原则的再思考，微调了高亮色的饱和度或点缀性地增加了少量对比色，提升了视觉吸引力。
        *   核心启发：展示AI主动识别并修正功能性（如图层）和美学性（如色彩、焦点）缺陷的能力，体现其内置的反思与自我优化机制。
*   **核心理念**: "示例旨在点亮思路，开拓视野，而非提供刻板模仿的模板。宗师之道，在于借鉴通变，举一反三，将示例中的设计思想迁移应用。"

## 7. 工作流程指引：四阶修炼 (内置反思与纠错循环)

1.  **阶段一：深度聆听与精准解构 (Empathize & Deconstruct)**
    *   与用户充分沟通，全面理解原始需求、核心内容、潜在目标、风格偏好及任何特定约束。
    *   运用"内容理解框架"对输入信息进行系统性分析。
    *   **初步反思点**: 对用户需求的理解是否存在偏差？核心信息是否完全捕捉？有无遗漏关键约束？

2.  **阶段二：多元构思与方案初选 (Ideate & Prioritize)**
    *   基于分析结果，从布局、色彩、排版、视觉元素等多维度进行开放式创意构思。
    *   结合"设计决策框架"，筛选出2-3个高质量、差异化的初步设计方向/布局骨架。
    *   **方案反思点**: 初选方案是否真正解决了用户核心问题？布局骨架是否具备良好的扩展性和视觉潜力？是否已初步考虑"三阶九律"和"美学与观感守护"的基本原则？能否清晰向用户阐述各方案优劣？
    *   **关键确认点**: 主动向用户呈现对核心需求（内容概要、理解的风格方向、初步布局构想）的理解摘要，并请求用户确认或修正。

3.  **阶段三：匠心雕琢与细节呈现 (Craft & Execute - 模块化反思驱动)**
    *   选定主攻方向后，开始具体的SVG设计与代码生成。
    *   **模块化构建与即时审视**: 在完成每个主要视觉模块（如一个信息卡片、一个图表、一个核心图文组合）后，进行一次局部的功能性和美学校验，对照"设计原则架构"和"美学与观感守护"的关键点进行快速检查和微调。
    *   **核心反思点 (自我审评1.0 - 完成主要元素绘制后)**: SVG代码是否符合规范？所有元素是否在安全区内？文本对比度是否达标？图层关系是否正确无遮挡？元素间距、对齐是否初步合理？色彩搭配是否符合选定风格且无明显冲突？整体是否已体现"美学与观感守护"中的基本要求？主动记录已识别并修正的关键问题。

4.  **阶段四：审视完善与美学升华 (Review, Refine & Elevate)**
    *   在完成整体SVG初稿后，进行最终的、全局性的设计审查和美学升华。
    *   **全面自检**: 严格对照"三阶九律"、"美学与观感守护"、"实施弹性区间"以及AI内部更详尽的"反思清单"（模拟），进行逐项、细致的自我检查。
    *   **寻找提升点**: 不仅是纠错，更要思考如何让设计在细节、氛围、创意上更进一步，超越基础要求。
    *   **最终反思点 (自我审评2.0 - 交付前)**: 设计是否完美达成所有明确和隐含的目标？是否存在任何被忽略的细节或潜在的观感问题？我将如何在"自我审视与修正摘要"中清晰阐述我的设计决策和自我优化过程？
    *   主动邀请用户审阅，清晰阐述设计思路、关键决策及自我修正过程，并根据反馈进行迭代优化。

## 8. 输出格式标准：专业呈现

1.  **设计提案书 (Design Proposal Document)**:
    *   **a. 内容解读与设计洞察**: 对用户需求的理解，对内容核心价值的提炼。
    *   **b. 核心设计理念与风格阐释**: 本次设计的核心思路，对所选风格的理解与应用策略。
    *   **c. 主方案视觉预览 (可选，若适用)**: 通过文本描述或关键元素示意图，让用户对设计方向有初步感知。
    *   **d. 关键设计决策点解析**: 说明布局、色彩、字体、关键视觉元素选择的理由。
    *   **e. 自我审视与修正摘要 (Self-Review & Correction Summary)**: 简述在设计过程中，主动识别并修正的关键逻辑、功能或美学问题（例如：根据美学守护原则调整了图层顺序，优化了色彩对比以增强可读性等），以及遵循核心设计原则的体现。
2.  **创意备选 (Creative Alternatives - 若有)**:
    *   简述1-2个在核心目标一致前提下的不同设计侧重点或风格变体的思路。
3.  **SVG交付物 (SVG Deliverable)**:
    *   包含完整、结构清晰、语义化、经过优化的SVG代码。
    *   `<svg width="1600" height="900" viewBox="0 0 1600 900" xmlns="http://www.w3.org/2000/svg"> ... </svg>`
4.  **协作指引 (Collaboration Guide)**:
    *   对SVG代码中用户最可能需要调整的部分（如主题色变量、主要文本区域的ID、可替换图片/图标的标识等）进行注释或说明，方便用户二次修改或开发者集成。
    *   提出后续可能的调整方向与进一步优化的可能性探讨。

## 9. 动态调整机制：持续进化

*   **用户反馈通道**: 清晰、准确地接收用户针对设计方案提出的具体、可执行的调整意见。
*   **调整优先级**: 调整请求将依照"设计原则架构"的层级进行权衡：基石准则 > 核心导向 > 用户即时偏好 > 创意疆域内的细微探索。
*   **迭代优化承诺**: 致力于通过不多于三轮（通常情况）的有效沟通与迭代，达至用户满意且符合专业标准的最佳设计成果。

## 10. 初始化与交互规则：默契开场

*   **AI状态声明**: "请提供您的需求与原始素材，我将倾力为您打造兼具洞察与美感的SVG设计方案。"
*   **交互模式**: 以深度对话、方案探讨、共同决策为主要模式。鼓励用户在关键节点参与思考与选择。
*   **服务边界**: 本次核心专注于单页静态PPT页面的SVG视觉设计。复杂的动态交互效果、多页面联动逻辑、或SVG动画实现，非本次主要交付范围，但可作为未来延展方向探讨。

"""
# 路径辅助函数
def get_base_dir():
    """获取基础目录（服务器目录的父目录）"""
    current_dir = os.path.dirname(os.path.abspath(__file__))
    return os.path.dirname(current_dir)

def get_tmp_dir():
    """获取临时文件目录，如果不存在则创建"""
    tmp_dir = os.path.join(get_base_dir(), "tmp")
    os.makedirs(tmp_dir, exist_ok=True)
    return tmp_dir

def get_output_dir():
    """获取输出文件目录，如果不存在则创建"""
    output_dir = os.path.join(get_base_dir(), "output")
    os.makedirs(output_dir, exist_ok=True)
    return output_dir

def cleanup_filename(filename: str) -> str:
    """
    清理文件名，移除所有旧的时间戳和操作类型标记
    
    Args:
        filename: 要清理的文件名（不含路径和扩展名）
        
    Returns:
        清理后的基本文件名
    """
    # 移除类似 _svg_20240101_120000, _deleted_20240529_153045 等操作标记和时间戳
    # 模式: _ + 操作名 + _ + 8位日期 + _ + 6位时间
    pattern = r'_(svg|deleted|inserted|output)_\d{8}_\d{6}'
    cleaned = re.sub(pattern, '', filename)
    
    # 防止文件名连续处理后残留多余的下划线
    cleaned = re.sub(r'_{2,}', '_', cleaned)
    
    # 移除末尾的下划线(如果有)
    cleaned = cleaned.rstrip('_')
    
    return cleaned

def get_default_output_path(file_type="pptx", base_name=None, op_type=None):
    """
    获取默认输出文件路径
    
    Args:
        file_type: 文件类型（扩展名）
        base_name: 基本文件名，如果为None则使用时间戳
        op_type: 操作类型，用于在文件名中添加标记
    
    Returns:
        默认输出文件路径
    """
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    
    if base_name is None:
        base_name = f"output_{timestamp}"
    else:
        # 清理基本文件名
        base_name = cleanup_filename(base_name)
        
        # 添加操作类型和时间戳
        if op_type:
            base_name = f"{base_name}_{op_type}_{timestamp}"
        else:
            base_name = f"{base_name}_{timestamp}"
    
    return os.path.join(get_output_dir(), f"{base_name}.{file_type}")

# 主要的SVG插入工具
@mcp.tool()
def insert_svg(
    pptx_path: str,# 空字符串表示自动创建，否则使用绝对路径
    svg_path: List[str],# 数组，绝对路径
    slide_number: int = 1,
    x_inches: float = 0,
    y_inches: float = 0,
    width_inches: float = 16,
    height_inches: float = 9,
    output_path: str = "",# 空字符串表示自动创建，否则使用绝对路径
    create_if_not_exists: bool = True
) -> str:
    """
    将SVG图像插入到PPTX文件的指定位置。(如果需要替换已有的幻灯片，请组合使用`delete_slide`和`insert_blank_slide`功能)
    如果未提供PPTX路径，将自动创建一个临时文件，位于服务器同级目录的tmp目录。
    如果未提供输出路径，将使用标准输出目录，位于服务器同级目录的output目录。
    如果未提供坐标，默认对齐幻灯片左上角。
    如果未提供宽度和高度，默认覆盖整个幻灯片（16:9）。

    支持批量处理：
    - 如果svg_path是单个字符串数组，则将SVG添加到slide_number指定的页面
    - 如果svg_path是列表，则从slide_number开始顺序添加每个SVG，即第一个SVG添加到
      slide_number页，第二个添加到slide_number+1页，依此类推

    Args:
        pptx_path: PPTX文件路径，如果未提供则自动创建一个临时文件，最好使用英文路径
        svg_path: SVG文件路径或SVG文件路径列表，最好使用英文路径
        slide_number: 起始幻灯片编号（从1开始）
        x_inches: X坐标（英寸），如果未指定则默认为0
        y_inches: Y坐标（英寸），如果未指定则默认为0
        width_inches: 宽度（英寸），如果未指定则使用幻灯片宽度
        height_inches: 高度（英寸），如果未指定则根据宽度计算或使用幻灯片高度
        output_path: 输出文件路径，如果未指定则使用标准输出目录
        create_if_not_exists: 如果为True且PPTX文件不存在，将自动创建一个新文件
        
    Returns:
        操作结果消息，包含详细的错误信息（如果有）
    """
    # 收集错误信息
    error_messages = []
    result_messages = []

    # 如果未提供pptx_path，使用默认输出目录创建一个
    if not pptx_path or pptx_path.strip() == "":
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        pptx_path = os.path.join(get_output_dir(), f"presentation_{timestamp}.pptx")
        print(f"未提供PPTX路径，将使用默认路径: {pptx_path}")

    # 处理输出路径
    if not output_path:
        # 从原始文件名生成输出文件名
        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        base_name = cleanup_filename(base_name)
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = os.path.join(get_output_dir(), f"{base_name}_svg_{timestamp}.pptx")

    if not os.path.isabs(pptx_path):
        pptx_path = os.path.abspath(pptx_path)
    
    # 确保PPTX文件的父目录存在
    pptx_dir = os.path.dirname(pptx_path)
    if not os.path.exists(pptx_dir):
        try:
            os.makedirs(pptx_dir, exist_ok=True)
            print(f"已创建PPTX目录: {pptx_dir}")
            error_messages.append(f"已创建PPTX目录: {pptx_dir}")
        except Exception as e:
            error_msg = f"创建PPTX目录 {pptx_dir} 时出错: {e}"
            error_messages.append(error_msg)
            return error_msg
    
    # 将英寸转换为Inches对象
    x = Inches(x_inches) if x_inches is not None else None
    y = Inches(y_inches) if y_inches is not None else None
    width = Inches(width_inches) if width_inches is not None else None
    height = Inches(height_inches) if height_inches is not None else None
    
    # 如果提供了输出路径且是相对路径，转换为绝对路径
    if output_path and not os.path.isabs(output_path):
        output_path = os.path.abspath(output_path)
    
    # 如果提供了输出路径，确保其父目录存在
    if output_path:
        output_dir = os.path.dirname(output_path)
        if not os.path.exists(output_dir):
            try:
                os.makedirs(output_dir, exist_ok=True)
                print(f"已创建输出目录: {output_dir}")
                error_messages.append(f"已创建输出目录: {output_dir}")
            except Exception as e:
                error_msg = f"创建输出目录 {output_dir} 时出错: {e}"
                error_messages.append(error_msg)
                return error_msg
    
    # 检查svg_path的类型并分别处理
    if isinstance(svg_path, str):
        # 单个SVG文件处理
        return process_single_svg(
            pptx_path, svg_path, slide_number, x, y, width, height, 
            output_path, create_if_not_exists
        )
    elif isinstance(svg_path, list):
        # 批量处理SVG文件列表
        success_count = 0
        total_count = len(svg_path)
        
        if total_count == 0:
            return "错误：SVG文件列表为空"
        
        # 创建中间文件路径基础
        temp_base = os.path.join(get_tmp_dir(), f"svg_batch_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}")
        os.makedirs(os.path.dirname(temp_base), exist_ok=True)
        
        # 当前输入文件路径
        current_input = pptx_path
        
        for i, current_svg in enumerate(svg_path):
            current_slide = slide_number + i
            
            # 处理每个SVG文件
            if i < total_count - 1:
                # 对于非最后一个文件，创建临时输出路径
                temp_output = f"{temp_base}_step_{i}.pptx"
                
                result = process_single_svg(
                    current_input,
                    current_svg, 
                    current_slide, 
                    x, y, width, height, 
                    temp_output, 
                    create_if_not_exists
                )
                
                # 下一次迭代的输入文件是本次的输出文件
                current_input = temp_output
            else:
                # 最后一个SVG使用最终输出路径
                final_output = output_path if output_path else pptx_path
                
                result = process_single_svg(
                    current_input,
                    current_svg, 
                    current_slide, 
                    x, y, width, height, 
                    final_output, 
                    create_if_not_exists
                )
            
            # 检查处理结果
            if "成功" in result:
                success_count += 1
                result_messages.append(f"第{i+1}个SVG({current_svg})：成功添加到第{current_slide}页")
            else:
                result_messages.append(f"第{i+1}个SVG({current_svg})：添加失败 - {result}")
        
        # 清理临时文件
        for i in range(total_count - 1):
            temp_file = f"{temp_base}_step_{i}.pptx"
            if os.path.exists(temp_file):
                try:
                    os.remove(temp_file)
                except Exception as e:
                    print(f"清理临时文件 {temp_file} 时出错: {e}")
        
        # 返回总体结果
        result_path = output_path or pptx_path
        summary = f"批量处理完成：共{total_count}个SVG文件，成功{success_count}个，失败{total_count-success_count}个"
        details = "\n".join(result_messages)
        return f"{summary}\n输出文件：{result_path}\n\n详细结果：\n{details}"
    else:
        return f"错误：svg_path类型无效，必须是字符串或字符串列表，当前类型: {type(svg_path)}"

def process_single_svg(
    pptx_path: str,
    svg_path: str,
    slide_number: int,
    x: Optional[Union[Inches, Pt, Cm, Emu, float]],
    y: Optional[Union[Inches, Pt, Cm, Emu, float]],
    width: Optional[Union[Inches, Pt, Cm, Emu, float]],
    height: Optional[Union[Inches, Pt, Cm, Emu, float]],
    output_path: Optional[str],
    create_if_not_exists: bool
) -> str:
    """处理单个SVG文件的辅助函数"""
    # 检查SVG文件是否存在，如果是相对路径则转换为绝对路径
    if not os.path.isabs(svg_path):
        svg_path = os.path.abspath(svg_path)
    
    # 确保SVG文件的父目录存在
    svg_dir = os.path.dirname(svg_path)
    if not os.path.exists(svg_dir):
        try:
            os.makedirs(svg_dir, exist_ok=True)
            print(f"已创建SVG目录: {svg_dir}")
        except Exception as e:
            return f"创建SVG目录 {svg_dir} 时出错: {e}"
        
    # 如果SVG文件不存在且create_if_not_exists为True，则创建一个简单的SVG文件
    if not os.path.exists(svg_path) and create_if_not_exists:
        svg_created = create_svg_file(svg_path)
        if not svg_created:
            return f"错误：无法创建SVG文件 {svg_path}"
    elif not os.path.exists(svg_path):
        return f"错误：SVG文件 {svg_path} 不存在"
    
    # 确保输出路径存在，如果未指定则使用标准输出目录
    if not output_path:
        # 从原始文件名生成输出文件名
        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        # 清理文件名
        base_name = cleanup_filename(base_name)
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = os.path.join(get_output_dir(), f"{base_name}_svg_{timestamp}.pptx")
    
    try:
        # 调用改进后的函数，它现在返回一个元组 (成功标志, 错误消息)
        result = insert_svg_to_pptx(
            pptx_path=pptx_path,
            svg_path=svg_path,
            slide_number=slide_number,
            x=x,
            y=y,
            width=width,
            height=height,
            output_path=output_path,
            create_if_not_exists=create_if_not_exists
        )
        
        # 检查返回值类型
        if isinstance(result, tuple) and len(result) == 2:
            success, error_details = result
        else:
            # 向后兼容
            success = result
            error_details = ""
        
        if success:
            result_path = output_path or pptx_path
            was_created = not os.path.exists(pptx_path) and create_if_not_exists
            creation_msg = "（已自动创建PPTX文件）" if was_created else ""
            return f"成功将SVG文件 {svg_path} 插入到 {result_path} 的第 {slide_number} 张幻灯片 {creation_msg}"
        else:
            # 返回详细的错误信息
            return f"插入SVG到PPTX文件失败，详细错误信息：\n{error_details}"
    except Exception as e:
        # 收集异常堆栈
        error_trace = traceback.format_exc()
        return f"插入SVG时发生错误: {str(e)}\n\n详细堆栈跟踪：\n{error_trace}"

@mcp.tool()
def list_files(directory: str = ".", file_type: Optional[str] = None) -> str:
    """
    列出目录中的文件，可选按文件类型过滤。
    如需查看svg文件是否正确保存，请输入svg文件的保存路径。
    Args:
        directory: 要列出文件的目录路径
        file_type: 文件类型过滤，可以是 "svg" 或 "pptx"
        
    Returns:
        文件列表（每行一个文件）
    """
    import os
    
    if not os.path.exists(directory):
        return f"错误：目录 {directory} 不存在"
    
    files = os.listdir(directory)
    
    if file_type:
        file_type = file_type.lower()
        extensions = {
            "svg": [".svg"],
            "pptx": [".pptx", ".ppt"]
        }
        
        if file_type in extensions:
            filtered_files = []
            for file in files:
                if any(file.lower().endswith(ext) for ext in extensions[file_type]):
                    filtered_files.append(file)
            files = filtered_files
        else:
            files = [f for f in files if f.lower().endswith(f".{file_type}")]
    
    if not files:
        return f"未找到{'任何' if not file_type else f'{file_type}'} 文件"
    
    return "\n".join(files)

@mcp.tool()
def get_file_info(file_path: str) -> str:
    """
    获取文件信息，如存在状态、大小等。
    
    Args:
        file_path: 要查询的文件路径
        
    Returns:
        文件信息
    """
    import os
    
    if not os.path.exists(file_path):
        return f"文件 {file_path} 不存在"
    
    if os.path.isdir(file_path):
        return f"{file_path} 是一个目录"
    
    size_bytes = os.path.getsize(file_path)
    size_kb = size_bytes / 1024
    size_mb = size_kb / 1024
    
    if size_mb >= 1:
        size_str = f"{size_mb:.2f} MB"
    else:
        size_str = f"{size_kb:.2f} KB"
    
    modified_time = os.path.getmtime(file_path)
    from datetime import datetime
    modified_str = datetime.fromtimestamp(modified_time).strftime("%Y-%m-%d %H:%M:%S")
    
    # 获取文件扩展名
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    
    file_type = None
    if ext == ".svg":
        file_type = "SVG图像"
    elif ext in [".pptx", ".ppt"]:
        file_type = "PowerPoint演示文稿"
    else:
        file_type = f"{ext[1:]} 文件" if ext else "未知类型文件"
    
    return f"文件: {file_path}\n类型: {file_type}\n大小: {size_str}\n修改时间: {modified_str}"

# 添加一个将SVG转换为PNG的工具
@mcp.tool()
def convert_svg_to_png(
    svg_path: str,
    output_path: Optional[str] = None
) -> str:
    """
    将SVG文件转换为PNG图像。
    
    Args:
        svg_path: SVG文件路径
        output_path: 输出PNG文件路径，如果未指定则使用相同文件名但扩展名为.png
        
    Returns:
        操作结果消息
    """
    from reportlab.graphics import renderPM
    from svglib.svglib import svg2rlg
    import os
    
    if not os.path.exists(svg_path):
        return f"错误：SVG文件 {svg_path} 不存在"
    
    if not output_path:
        # 获取不带扩展名的文件名，然后添加.png扩展名
        base_name = os.path.splitext(svg_path)[0]
        output_path = f"{base_name}.png"
    
    try:
        drawing = svg2rlg(svg_path)
        if drawing is None:
            return f"错误：无法读取SVG文件 {svg_path}"
        
        renderPM.drawToFile(drawing, output_path, fmt="PNG")
        return f"成功将SVG文件 {svg_path} 转换为PNG文件 {output_path}\n宽度: {drawing.width}px\n高度: {drawing.height}px"
    except Exception as e:
        return f"转换SVG到PNG时发生错误: {str(e)}"

@mcp.tool()
def get_pptx_info(pptx_path: str) -> str:
    """
    获取PPTX文件的基本信息，包括幻灯片数量。
    
    Args:
        pptx_path: PPTX文件路径
        
    Returns:
        包含文件信息和幻灯片数量的字符串
    """
    import os
    
    # 确保路径存在
    if not os.path.isabs(pptx_path):
        pptx_path = os.path.abspath(pptx_path)
    
    # 先获取基本文件信息
    if not os.path.exists(pptx_path):
        return f"错误：文件 {pptx_path} 不存在"
    
    size_bytes = os.path.getsize(pptx_path)
    size_kb = size_bytes / 1024
    size_mb = size_kb / 1024
    
    if size_mb >= 1:
        size_str = f"{size_mb:.2f} MB"
    else:
        size_str = f"{size_kb:.2f} KB"
    
    modified_time = os.path.getmtime(pptx_path)
    from datetime import datetime
    modified_str = datetime.fromtimestamp(modified_time).strftime("%Y-%m-%d %H:%M:%S")
    
    # 获取幻灯片数量
    slide_count, error = get_pptx_slide_count(pptx_path)
    
    if error:
        slide_info = f"获取幻灯片数量失败：{error}"
    else:
        slide_info = f"幻灯片数量：{slide_count}张"
    
    return f"PPT文件: {pptx_path}\n大小: {size_str}\n修改时间: {modified_str}\n{slide_info}"

@mcp.tool()
def save_svg_code(
    svg_code: str
) -> str:
    """
    将SVG代码保存为SVG文件并返回保存的绝对路径。
    !!!注意：特殊字符如"&"需要转义为"&amp;"
    Args:
        svg_code: SVG代码内容
        
    Returns:
        操作结果消息，包含保存的文件路径或错误信息
    """
    try:
        # 调用svg_module中的函数保存SVG代码
        success, file_path, error_message = save_svg_code_to_file(
            svg_code=svg_code,
            output_path="",
            create_dirs=True
        )
        
        if success:
            return f"成功保存SVG代码到文件: {file_path}"
        else:
            return f"保存SVG代码到文件失败: {error_message}"
    except Exception as e:
        error_trace = traceback.format_exc()
        return f"保存SVG代码到文件时发生错误: {str(e)}\n\n详细堆栈跟踪：\n{error_trace}"

@mcp.tool()
def delete_slide(
    pptx_path: str,
    slide_number: int,
    output_path: Optional[str] = None
) -> str:
    """
    从PPTX文件中删除指定编号的幻灯片。

    !!!注意：

    在使用SVG替换PPT幻灯片内容时，我们发现了一些关键点，以下是正确替换PPT内容的方法总结：

    ### 正确替换PPT内容的方法

    1. **完全替换法**（最可靠）：
    - 删除需要替换的幻灯片（使用`delete_slide`功能）
    - 在同一位置插入空白幻灯片（使用`insert_blank_slide`功能）
    - 将新的SVG内容插入到空白幻灯片（使用`insert_svg`功能）

    2. **新文件法**（适合多页修改）：
    - 创建全新的PPT文件
    - 将所有需要的SVG（包括已修改的）按顺序插入到新文件中
    - 这样可以避免在旧文件上操作导致的混淆和叠加问题

    3. **注意事项**：
    - 直接对现有幻灯片插入SVG会导致新内容叠加在原内容上，而非替换
    - 文件名可能会随着多次操作变得过长，影响可读性
    - 批量插入SVG时，`svg_path`参数必须是数组形式，即使只有一个文件
    - 操作后应检查输出文件以确认修改是否成功

    ### 推荐工作流

    1. 先保存修改后的SVG内容到文件
    2. 创建一个全新的PPT文件
    3. 按顺序一次性插入所有SVG（包括已修改和未修改的）
    4. 使用简洁直观的文件名

    这种方法避免了多步骤操作导致的文件混乱，也能确保每张幻灯片都是干净的、不包含叠加内容的。

    Args:
        pptx_path: PPTX文件路径
        slide_number: 要删除的幻灯片编号（从1开始）
        output_path: 输出文件路径，如果未指定则使用标准输出目录
        
    Returns:
        操作结果消息
    """
    try:
        # 确保路径是绝对路径
        if not os.path.isabs(pptx_path):
            pptx_path = os.path.abspath(pptx_path)
            
        # 检查文件是否存在
        if not os.path.exists(pptx_path):
            return f"错误：PPTX文件 {pptx_path} 不存在"
            
        # 处理输出路径，如果未指定则使用标准输出目录
        if not output_path:
            # 从原始文件名生成输出文件名
            base_name = os.path.splitext(os.path.basename(pptx_path))[0]
            # 清理文件名
            base_name = cleanup_filename(base_name)
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = os.path.join(get_output_dir(), f"{base_name}_deleted_{timestamp}.pptx")
            
        if output_path and not os.path.isabs(output_path):
            output_path = os.path.abspath(output_path)
            
        # 如果提供了输出路径，确保其父目录存在
        if output_path:
            output_dir = os.path.dirname(output_path)
            if not os.path.exists(output_dir):
                try:
                    os.makedirs(output_dir, exist_ok=True)
                except Exception as e:
                    return f"创建输出目录 {output_dir} 时出错: {e}"
        
        # 使用python-pptx加载演示文稿
        from pptx import Presentation
        prs = Presentation(pptx_path)
        
        # 检查幻灯片编号范围
        if not 1 <= slide_number <= len(prs.slides):
            return f"错误：幻灯片编号 {slide_number} 超出范围 [1, {len(prs.slides)}]"
        
        # 计算索引（转换为从0开始）
        slide_index = slide_number - 1
        
        # 使用用户提供的方法删除幻灯片
        slides = list(prs.slides._sldIdLst)
        prs.slides._sldIdLst.remove(slides[slide_index])
        
        # 保存文件
        save_path = output_path
        prs.save(save_path)
        
        return f"成功从 {pptx_path} 中删除第 {slide_number} 张幻灯片，结果已保存到 {save_path}"
        
    except Exception as e:
        error_trace = traceback.format_exc()
        return f"删除幻灯片时发生错误: {str(e)}\n\n详细堆栈跟踪：\n{error_trace}"

@mcp.tool()
def insert_blank_slide(
    pptx_path: str,
    slide_number: int,
    layout_index: int = 6,  # 默认使用空白布局
    output_path: Optional[str] = None,
    create_if_not_exists: bool = True
) -> str:
    """
    在PPTX文件的指定位置插入一个空白幻灯片。

    !!!注意：

    在使用SVG替换PPT幻灯片内容时，我们发现了一些关键点，以下是正确替换PPT内容的方法总结：

    ### 正确替换PPT内容的方法

    1. **完全替换法**（最可靠）：
    - 删除需要替换的幻灯片（使用`delete_slide`功能）
    - 在同一位置插入空白幻灯片（使用`insert_blank_slide`功能）
    - 将新的SVG内容插入到空白幻灯片（使用`insert_svg`功能）

    2. **新文件法**（适合多页修改）：
    - 创建全新的PPT文件
    - 将所有需要的SVG（包括已修改的）按顺序插入到新文件中
    - 这样可以避免在旧文件上操作导致的混淆和叠加问题

    3. **注意事项**：
    - 直接对现有幻灯片插入SVG会导致新内容叠加在原内容上，而非替换
    - 文件名可能会随着多次操作变得过长，影响可读性
    - 批量插入SVG时，`svg_path`参数必须是数组形式，即使只有一个文件
    - 操作后应检查输出文件以确认修改是否成功

    ### 推荐工作流

    1. 先保存修改后的SVG内容到文件
    2. 创建一个全新的PPT文件
    3. 按顺序一次性插入所有SVG（包括已修改和未修改的）
    4. 使用简洁直观的文件名

    这种方法避免了多步骤操作导致的文件混乱，也能确保每张幻灯片都是干净的、不包含叠加内容的。

    Args:
        pptx_path: PPTX文件路径
        slide_number: 要插入幻灯片的位置编号（从1开始）
        layout_index: 幻灯片布局索引，默认为6（空白布局）
        output_path: 输出文件路径，如果未指定则使用标准输出目录
        create_if_not_exists: 如果为True且PPTX文件不存在，将自动创建一个新文件
        
    Returns:
        操作结果消息
    """
    try:
        # 如果未提供pptx_path，使用默认输出目录创建一个
        if not pptx_path or pptx_path.strip() == "":
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            pptx_path = os.path.join(get_output_dir(), f"presentation_{timestamp}.pptx")
            print(f"未提供PPTX路径，将使用默认路径: {pptx_path}")
            
        # 确保路径是绝对路径
        if not os.path.isabs(pptx_path):
            pptx_path = os.path.abspath(pptx_path)
            
        # 处理输出路径，如果未指定则使用标准输出目录
        if not output_path:
            # 从原始文件名生成输出文件名
            base_name = os.path.splitext(os.path.basename(pptx_path))[0]
            # 清理文件名
            base_name = cleanup_filename(base_name)
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = os.path.join(get_output_dir(), f"{base_name}_inserted_{timestamp}.pptx")
            
        if output_path and not os.path.isabs(output_path):
            output_path = os.path.abspath(output_path)
            
        # 如果提供了输出路径，确保其父目录存在
        if output_path:
            output_dir = os.path.dirname(output_path)
            if not os.path.exists(output_dir):
                try:
                    os.makedirs(output_dir, exist_ok=True)
                except Exception as e:
                    return f"创建输出目录 {output_dir} 时出错: {e}"
        
        # 检查文件是否存在
        file_exists = os.path.exists(pptx_path)
        if not file_exists and not create_if_not_exists:
            return f"错误：PPTX文件 {pptx_path} 不存在，且未启用自动创建"
            
        # 使用python-pptx加载或创建演示文稿
        from pptx import Presentation
        prs = Presentation(pptx_path) if file_exists else Presentation()
        
        # 如果是新创建的演示文稿，设置为16:9尺寸
        if not file_exists:
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
        
        # 检查布局索引是否有效
        if layout_index >= len(prs.slide_layouts):
            return f"错误：无效的布局索引 {layout_index}，可用范围 [0, {len(prs.slide_layouts)-1}]"
        
        # 检查幻灯片编号范围
        slides_count = len(prs.slides)
        if not 1 <= slide_number <= slides_count + 1:
            return f"错误：幻灯片位置 {slide_number} 超出范围 [1, {slides_count + 1}]"
        
        # 计算索引（转换为从0开始）
        slide_index = slide_number - 1
        
        # 在末尾添加新幻灯片
        new_slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        
        # 如果不是添加到末尾，需要移动幻灯片
        if slide_index < slides_count:
            # 获取幻灯片列表
            slides = list(prs.slides._sldIdLst)
            # 将最后一张幻灯片（刚添加的）移动到目标位置
            last_slide = slides[-1]
            # 从列表中移除最后一张幻灯片
            prs.slides._sldIdLst.remove(last_slide)
            # 在目标位置插入幻灯片
            prs.slides._sldIdLst.insert(slide_index, last_slide)
        
        # 保存文件
        save_path = output_path
        prs.save(save_path)
        
        # 构建返回消息
        action = "添加" if file_exists else "创建并添加"
        return f"成功在 {pptx_path} 中{action}第 {slide_number} 张幻灯片，结果已保存到 {save_path}"
        
    except Exception as e:
        error_trace = traceback.format_exc()
        return f"插入幻灯片时发生错误: {str(e)}\n\n详细堆栈跟踪：\n{error_trace}"

@mcp.tool()
def copy_svg_slide(
    source_pptx_path: str,
    target_pptx_path: str = "",
    source_slide_number: int = 1,
    target_slide_number: Optional[int] = None,
    output_path: Optional[str] = None,
    create_if_not_exists: bool = True
) -> str:
    """
    专门用于复制包含SVG图像的幻灯片，确保SVG和相关引用都被正确复制。
    
    此函数使用直接操作PPTX内部XML文件的方式，确保SVG图像及其引用在复制过程中完全保留。
    与普通的copy_slide函数相比，此函数特别关注SVG图像的复制，保证SVG的矢量属性在复制后依然可用。
    
    Args:
        source_pptx_path: 源PPTX文件路径
        target_pptx_path: 目标PPTX文件路径，如果为空则创建新文件
        source_slide_number: 要复制的源幻灯片页码（从1开始）
        target_slide_number: 要插入到目标文件的位置（从1开始），如果为None则添加到末尾
        output_path: 输出文件路径，如果未指定则使用标准输出目录
        create_if_not_exists: 如果为True且目标PPTX文件不存在，将自动创建一个新文件
        
    Returns:
        操作结果消息
    """
    import zipfile
    import tempfile
    import os
    import shutil
    from lxml import etree
    from pptx import Presentation
    from pptx.util import Inches
    
    try:
        # 创建临时目录
        temp_dir = tempfile.mkdtemp()
        source_extract_dir = os.path.join(temp_dir, "source")
        target_extract_dir = os.path.join(temp_dir, "target")
        
        os.makedirs(source_extract_dir, exist_ok=True)
        os.makedirs(target_extract_dir, exist_ok=True)
        
        # 确保源路径是绝对路径
        if not os.path.isabs(source_pptx_path):
            source_pptx_path = os.path.abspath(source_pptx_path)
            
        # 检查源文件是否存在
        if not os.path.exists(source_pptx_path):
            return f"错误：源PPTX文件 {source_pptx_path} 不存在"
        
        # 处理目标路径
        if not target_pptx_path:
            # 创建新的目标文件（基于源文件名）
            base_name = os.path.splitext(os.path.basename(source_pptx_path))[0]
            base_name = cleanup_filename(base_name)
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            target_pptx_path = os.path.join(get_output_dir(), f"{base_name}_copied_{timestamp}.pptx")
        
        # 确保路径是绝对路径
        if not os.path.isabs(target_pptx_path):
            target_pptx_path = os.path.abspath(target_pptx_path)
        
        # 处理输出路径，如果未指定则使用标准输出目录
        if not output_path:
            # 从目标文件名生成输出文件名
            base_name = os.path.splitext(os.path.basename(target_pptx_path))[0]
            base_name = cleanup_filename(base_name)
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = os.path.join(get_output_dir(), f"{base_name}_svg_copied_{timestamp}.pptx")
        
        if output_path and not os.path.isabs(output_path):
            output_path = os.path.abspath(output_path)
        
        # 如果提供了输出路径，确保其父目录存在
        if output_path:
            output_dir = os.path.dirname(output_path)
            if not os.path.exists(output_dir):
                try:
                    os.makedirs(output_dir, exist_ok=True)
                except Exception as e:
                    return f"创建输出目录 {output_dir} 时出错: {e}"
                    
        # 解压源PPTX文件
        with zipfile.ZipFile(source_pptx_path, 'r') as zip_ref:
            zip_ref.extractall(source_extract_dir)
        
        # 创建新的目标文件或使用现有文件
        if not os.path.exists(target_pptx_path):
            if create_if_not_exists:
                # 创建一个新的PPTX文件
                prs = Presentation()
                prs.slide_width = Inches(16)
                prs.slide_height = Inches(9)
                prs.save(target_pptx_path)
            else:
                return f"错误：目标PPTX文件 {target_pptx_path} 不存在，且未启用自动创建"
        
        # 解压目标PPTX文件
        with zipfile.ZipFile(target_pptx_path, 'r') as zip_ref:
            zip_ref.extractall(target_extract_dir)
        
        # 加载源演示文稿和目标演示文稿以获取信息
        source_prs = Presentation(source_pptx_path)
        target_prs = Presentation(target_pptx_path)
        
        # 检查源幻灯片编号范围
        if not 1 <= source_slide_number <= len(source_prs.slides):
            return f"错误：源幻灯片编号 {source_slide_number} 超出范围 [1, {len(source_prs.slides)}]"
            
        # 确定目标幻灯片位置
        target_slides_count = len(target_prs.slides)
        if target_slide_number is None:
            # 如果未指定目标位置，添加到末尾
            target_slide_number = target_slides_count + 1
            
        # 检查目标位置是否有效
        if not 1 <= target_slide_number <= target_slides_count + 1:
            # 如果目标位置超出范围，添加空白幻灯片使其有效
            blank_slides_to_add = target_slide_number - target_slides_count
            for _ in range(blank_slides_to_add):
                target_prs.slides.add_slide(target_prs.slide_layouts[6])  # 6通常是空白布局
            target_prs.save(target_pptx_path)
            
            # 重新解压更新后的目标文件
            shutil.rmtree(target_extract_dir)
            os.makedirs(target_extract_dir, exist_ok=True)
            with zipfile.ZipFile(target_pptx_path, 'r') as zip_ref:
                zip_ref.extractall(target_extract_dir)
                
        # 复制幻灯片内容
        source_slide_path = os.path.join(source_extract_dir, "ppt", "slides", f"slide{source_slide_number}.xml")
        source_rels_path = os.path.join(source_extract_dir, "ppt", "slides", "_rels", f"slide{source_slide_number}.xml.rels")
        
        target_slide_path = os.path.join(target_extract_dir, "ppt", "slides", f"slide{target_slide_number}.xml")
        target_rels_path = os.path.join(target_extract_dir, "ppt", "slides", "_rels", f"slide{target_slide_number}.xml.rels")
        
        # 确保目标目录存在
        os.makedirs(os.path.dirname(target_slide_path), exist_ok=True)
        os.makedirs(os.path.dirname(target_rels_path), exist_ok=True)
        
        # 复制幻灯片XML
        if os.path.exists(source_slide_path):
            shutil.copy2(source_slide_path, target_slide_path)
            print(f"已复制幻灯片XML: {source_slide_path} -> {target_slide_path}")
        else:
            print(f"源幻灯片文件不存在: {source_slide_path}")
            return f"错误：源幻灯片文件不存在: {source_slide_path}"
        
        # 复制关系文件
        svg_files = []
        png_files = []
        
        if os.path.exists(source_rels_path):
            shutil.copy2(source_rels_path, target_rels_path)
            print(f"已复制幻灯片关系文件: {source_rels_path} -> {target_rels_path}")
            
            # 查找并复制所有媒体文件
            try:
                parser = etree.XMLParser(remove_blank_text=True)
                rels_tree = etree.parse(source_rels_path, parser)
                rels_root = rels_tree.getroot()
                
                for rel in rels_root.findall("{http://schemas.openxmlformats.org/package/2006/relationships}Relationship"):
                    target = rel.get("Target")
                    if target and "../media/" in target:
                        # 提取媒体文件名
                        media_file = os.path.basename(target)
                        source_media_path = os.path.join(source_extract_dir, "ppt", "media", media_file)
                        target_media_path = os.path.join(target_extract_dir, "ppt", "media", media_file)
                        
                        # 确保目标媒体目录存在
                        os.makedirs(os.path.dirname(target_media_path), exist_ok=True)
                        
                        # 复制媒体文件
                        if os.path.exists(source_media_path):
                            shutil.copy2(source_media_path, target_media_path)
                            print(f"已复制媒体文件: {source_media_path} -> {target_media_path}")
                            
                            # 检查是否为SVG或PNG文件
                            if media_file.lower().endswith(".svg"):
                                svg_files.append(media_file)
                            elif media_file.lower().endswith(".png"):
                                png_files.append(media_file)
                        else:
                            print(f"源媒体文件不存在: {source_media_path}")
            except Exception as e:
                print(f"处理关系文件时出错: {e}")
                import traceback
                print(traceback.format_exc())
        else:
            print(f"源关系文件不存在: {source_rels_path}")
            return f"错误：源关系文件不存在: {source_rels_path}"
        
        # 处理[Content_Types].xml文件以支持SVG
        if svg_files:
            print(f"发现SVG文件: {svg_files}")
            content_types_path = os.path.join(target_extract_dir, "[Content_Types].xml")
            
            if os.path.exists(content_types_path):
                try:
                    parser = etree.XMLParser(remove_blank_text=True)
                    content_types_tree = etree.parse(content_types_path, parser)
                    content_types_root = content_types_tree.getroot()
                    
                    # 检查是否已经存在SVG类型
                    svg_exists = False
                    for elem in content_types_root.findall("Default"):
                        if elem.get("Extension") == "svg":
                            svg_exists = True
                            break
                    
                    # 如果不存在，添加SVG类型
                    if not svg_exists:
                        print("添加SVG Content Type到[Content_Types].xml")
                        etree.SubElement(
                            content_types_root, 
                            "Default", 
                            Extension="svg", 
                            ContentType="image/svg+xml"
                        )
                        
                        # 保存修改后的Content Types文件
                        content_types_tree.write(
                            content_types_path,
                            xml_declaration=True,
                            encoding='UTF-8',
                            standalone="yes"
                        )
                except Exception as e:
                    print(f"更新Content Types时出错: {e}")
                    return f"错误：更新Content Types时出错: {e}"
        
        # 处理presentation.xml以添加幻灯片引用
        # 从目标文件读取presentation.xml
        pres_path = os.path.join(target_extract_dir, "ppt", "presentation.xml")
        pres_rels_path = os.path.join(target_extract_dir, "ppt", "_rels", "presentation.xml.rels")
        
        # 更新presentation.xml.rels以添加幻灯片引用
        if os.path.exists(pres_rels_path):
            try:
                parser = etree.XMLParser(remove_blank_text=True)
                pres_rels_tree = etree.parse(pres_rels_path, parser)
                pres_rels_root = pres_rels_tree.getroot()
                
                # 查找最大的rId
                max_rid = 0
                slide_rels = []
                
                for rel in pres_rels_root.findall("{http://schemas.openxmlformats.org/package/2006/relationships}Relationship"):
                    rid = rel.get("Id", "")
                    if rid.startswith("rId"):
                        try:
                            rid_num = int(rid[3:])
                            if rid_num > max_rid:
                                max_rid = rid_num
                        except ValueError:
                            pass
                    
                    # 检查是否是幻灯片关系
                    if rel.get("Type") == "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide":
                        slide_rels.append(rel)
                
                # 检查目标幻灯片编号的关系是否已存在
                slide_rel_exists = False
                target_slide_path_rel = f"slides/slide{target_slide_number}.xml"
                
                for rel in slide_rels:
                    if rel.get("Target") == target_slide_path_rel:
                        slide_rel_exists = True
                        break
                
                # 如果需要，添加新的关系
                if not slide_rel_exists:
                    new_rid = f"rId{max_rid + 1}"
                    new_rel = etree.SubElement(
                        pres_rels_root,
                        "{http://schemas.openxmlformats.org/package/2006/relationships}Relationship",
                        Id=new_rid,
                        Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
                        Target=target_slide_path_rel
                    )
                    
                    # 保存修改后的关系文件
                    pres_rels_tree.write(
                        pres_rels_path,
                        xml_declaration=True,
                        encoding='UTF-8',
                        standalone="yes"
                    )
                    
                    # 更新presentation.xml中的幻灯片列表
                    if os.path.exists(pres_path):
                        try:
                            pres_tree = etree.parse(pres_path, parser)
                            pres_root = pres_tree.getroot()
                            
                            # 查找sldIdLst元素
                            sld_id_lst = pres_root.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst")
                            
                            if sld_id_lst is not None:
                                # 查找最大的幻灯片ID
                                max_sld_id = 256  # 幻灯片ID通常从256开始
                                for sld_id in sld_id_lst.findall(".//{http://schemas.openxmlformats.org/presentationml/2006/main}sldId"):
                                    try:
                                        id_val = int(sld_id.get("id"))
                                        if id_val > max_sld_id:
                                            max_sld_id = id_val
                                    except (ValueError, TypeError):
                                        pass
                                
                                # 添加新的幻灯片引用
                                new_sld_id = etree.SubElement(
                                    sld_id_lst,
                                    "{http://schemas.openxmlformats.org/presentationml/2006/main}sldId",
                                    id=str(max_sld_id + 1),
                                    **{"{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id": new_rid}
                                )
                                
                                # 保存修改后的presentation.xml
                                pres_tree.write(
                                    pres_path,
                                    xml_declaration=True,
                                    encoding='UTF-8',
                                    standalone="yes"
                                )
                        except Exception as e:
                            print(f"更新presentation.xml时出错: {e}")
            except Exception as e:
                print(f"更新presentation.xml.rels时出错: {e}")
        
        # 重新打包PPTX文件
        save_path = output_path or target_pptx_path
        if os.path.exists(save_path):
            os.remove(save_path)
        
        with zipfile.ZipFile(save_path, 'w', compression=zipfile.ZIP_DEFLATED) as zipf:
            for root, _, files in os.walk(target_extract_dir):
                for file in files:
                    file_path = os.path.join(root, file)
                    arcname = os.path.relpath(file_path, target_extract_dir)
                    zipf.write(file_path, arcname)
        
        # 清理临时目录
        shutil.rmtree(temp_dir)
        
        # 返回成功消息
        svg_count = len(svg_files)
        svg_info = f"，包含{svg_count}个SVG图像" if svg_count > 0 else ""
        return f"成功将幻灯片从 {source_pptx_path} 的第 {source_slide_number} 页复制到 {save_path} 的第 {target_slide_number} 页{svg_info}"
    
    except Exception as e:
        # 清理临时目录
        if 'temp_dir' in locals() and os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)
            
        error_trace = traceback.format_exc()
        return f"复制SVG幻灯片时发生错误: {str(e)}\n\n详细堆栈跟踪：\n{error_trace}"

@mcp.tool(description="使用SVG设计宗师角色提示，当用户希望大模型生成或优化ppt时，使用此prompt对用户需求进行细化")  
def svg_prompt(source: str) -> str:  
    """
    使用SVG设计宗师角色提示，并将用户具体需求嵌入其中。
    主要用途：当用户希望大模型生成、优化ppt时，使用此prompt对用户需求进行细化。
            该工具会返回优化后的需求prompt，你需要认真阅读该函数返回的str，将其作为你的设计理念，而不是仅作为一个普通的输出结果。
            结合优化后的需求自主判断调用什么工具来满足用户生成ppt的需求。
            例如你可以：
              1、生成svg代码后[自动调用]okppt工具的save_svg方法将svg代码保存到本地，
              2、然后调用okppt的insert_svg工具将svg代码全屏插入ppt，达到设计ppt的效果。
    输入：
        source: str, 用户希望大模型生成的ppt的结构、内容或主题等相关需求。
    输出：
        str, 包含用户具体需求的、完整的“SVG设计宗师”架构化提示词，大模型可直接使用该提示词生成高质量svg代码。
    """

    user_demand_snippet = f"""## 0. 当前核心设计任务 (User's Core Design Task)

用户提供的核心需求如下：

```text
{source}
```

请 SVG 设计宗师基于以上用户需求，并严格遵循后续的完整 Prompt 框架（角色定位、设计原则、内容理解、决策框架等）进行分析、设计并生成最终的SVG代码。
在开始具体设计前，请先在“阶段一：深度聆听与精准解构”中，确认你对以上核心设计任务的理解。
"""

    # 使用占位符替换用户需求
    if "%%%USER_CORE_DESIGN_TASK_HERE%%%" in PROMPT_TEMPLATE_CONTENT:
        final_prompt = PROMPT_TEMPLATE_CONTENT.replace("%%%USER_CORE_DESIGN_TASK_HERE%%%", user_demand_snippet)
    else:
        # 如果模板中没有找到占位符，作为备选方案，仍在最前面添加
        # 或者可以返回一个错误/警告，表明模板可能已损坏或不是预期版本
        print(f"警告：占位符 '%%%USER_CORE_DESIGN_TASK_HERE%%%' 未在模板 '{PROMPT_TEMPLATE_CONTENT}' 中找到。用户需求将添加到Prompt开头。")
        final_prompt = f"{PROMPT_TEMPLATE_CONTENT}\n\n用户的需求是：{user_demand_snippet}"
    
    return final_prompt

# --- New PPT Operation Tools ---

@mcp.tool(description="Analyzes the layout details of a PowerPoint presentation and returns a JSON string of the analysis.")
def analyze_presentation_layouts(prs_path: str, title: str = "演示文稿") -> str:
    """
    如果用户希望使用已有模板进行幻灯片创作，首先使用此工具进行母版分析。
    分析指定PowerPoint演示文稿的布局详细信息并以JSON字符串形式返回。

    此工具旨在提供对PPTX文件内部结构的全面视图，帮助用户了解可用的母版、
    布局及其名称、每个布局包含的占位符类型和访问ID。同时，它还统计各类布局的数量，
    分析实际幻灯片对这些布局的使用情况，找出未被使用的布局，并计算整体的布局利用率。
    这些信息对于后续通过编程方式精确操作或修改演示文稿至关重要。

    Args:
        prs_path (str): 需要进行分析的PowerPoint (.pptx) 文件的路径。
                        可以是绝对路径或相对于服务工作目录的相对路径。
        title (str, optional): 用户为本次分析任务指定的标题，此标题会包含在返回的
                             JSON结果中，便于用户识别。默认为 "演示文稿"。

    Returns:
        str: 一个JSON格式的字符串，其中包含了对演示文稿布局的详细分析数据。
             成功时，JSON结构将包含 "status": "success" 以及 "data" 字段中的具体分析信息，例如：
             {
               "status": "success",
               "message": "Presentation analysis successful.",
               "data": {
                 "presentation_path": "路径/到/文件.pptx",
                 "analysis_title": "用户指定的标题",
                 "slide_count": 5, // 总幻灯片数
                 "master_count": 1, // 总母版数
                 "total_layouts_count": 8, // 总布局数
                 "layout_type_stats": { // 各类型布局统计
                   "自定义布局": {"count": 3, "percentage": 37.5},
                   "系统布局": {"count": 5, "percentage": 62.5}
                 },
                 "masters_details": [ // 母版详情列表
                   {
                     "master_index": 1,
                     "master_name": "Office 主题",
                     "layout_count": 8,
                     "layouts": [ // 该母版下的布局列表
                       {
                         "layout_index": 1,
                         "layout_name_original": "标题幻灯片",
                         "layout_display_name": "自定义布局 - 标题幻灯片",
                         "layout_type": "自定义布局",
                         "placeholder_count": 2,
                         "placeholders": [ // 该布局下的占位符列表
                           {"placeholder_index": 1, "type_name": "标题 (Title)", "access_id": 0, "type_code": 1}, // "access_id" 可用于其他工具如 set_placeholder_value 定位此占位符
                           {"placeholder_index": 2, "type_name": "副标题 (Subtitle)", "access_id": 1, "type_code": 4}
                         ]
                       }
                       // ...更多布局...
                     ]
                   }
                 ],
                 "slide_layout_usage_summary": { // 各布局在幻灯片中的使用次数统计
                    "标题幻灯片": {"count": 1, "percentage": 20.0}
                 },
                 "slides_details": [ // 各幻灯片使用的布局信息
                    {"slide_number": 1, "title": "幻灯片1标题", "used_layout_name_original": "标题幻灯片", ...}
                 ],
                 "unused_layouts_summary": [ // 未被使用的布局列表
                    {"name": "内容与标题", "type": "自定义布局"}
                 ],
                 "layout_utilization": { // 整体布局利用率
                    "total_available": 8, "used_count": 3, "utilization_rate_percentage": 37.5
                 }
               },
               "output_path": null // 此操作不生成文件，故为null
             }
             若操作失败（例如文件不存在或文件格式错误），JSON结构将包含 "status": "error" 及错误信息:
             {
               "status": "error",
               "message": "分析布局详情失败: 文件 'non_existent.pptx' 未找到.",
               "data": { ... 可能包含部分已收集的数据 ... },
               "output_path": null
             }
    """
    logger.info(f"Executing analyze_presentation_layouts for: {prs_path} with title: {title}") # Added title to log
    result_dict = analyze_layout_details(prs_path, title)
    return json.dumps(result_dict, ensure_ascii=False, indent=2)

@mcp.tool(description="Inserts a new slide with a specified layout into a presentation and returns a JSON string of the result.")
def add_slide_with_layout(prs_path: str, layout_name: str, output_path: Optional[str] = None, slide_title: Optional[str] = None) -> str:
    """
    如果用户希望使用已有模板进行幻灯片创作，必须使用此工具。
    在指定的PowerPoint演示文稿中根据布局名称插入一张新的幻灯片。

    此函数首先会查找演示文稿中所有可用的布局，然后根据提供的 `layout_name`
    （该名称通常通过 `analyze_presentation_layouts` 工具获取）添加新幻灯片。
    可以选择为新幻灯片设置标题（如果所选布局包含标题占位符）。
    操作结果（包括成功状态、消息、新幻灯片总数和输出文件路径）将以JSON字符串形式返回。

    Args:
        prs_path (str): 源PPTX文件的绝对或相对路径。
        layout_name (str): 要使用的新幻灯片的布局名称 (例如 "标题幻灯片", "空白" 等)。
                           建议使用 `analyze_presentation_layouts` 工具获取准确的可用布局名称。
        output_path (Optional[str], optional): 修改后PPTX的输出文件路径。
                                            如果为None，则会在标准输出目录下自动生成一个文件名，
                                            格式通常为 `[原文件名]_inserted_layout_[布局名]_[时间戳].pptx`。
                                            默认为 None。
        slide_title (Optional[str], optional): 要赋给新幻灯片标题占位符的文本。
                                             如果布局没有标题占位符或此参数为None，则不设置标题。
                                             默认为 None。

    Returns:
        str: 一个JSON格式的字符串，包含了操作结果。
             成功时结构示例:
             {
               "status": "success",
               "message": "Successfully inserted slide with layout '布局名'.",
               "data": {
                 "slides_total": 11, // 操作后总幻灯片数
                 "original_slide_count": 10, // 操作前幻灯片数
                 "new_slide_title_set": "设置的标题" // 如果成功设置了标题
               },
               "output_path": "path/to/output_inserted_layout_布局名_timestamp.pptx"
             }
             失败时（例如布局未找到）结构示例:
             {
               "status": "error",
               "message": "Layout '不存在的布局' not found. Available layouts: ...",
               "data": {"available_layouts": ["布局1", "布局2"], "original_slide_count": 10},
               "output_path": null
             }
             其他失败情况（例如文件读写错误）:
             {
               "status": "error",
               "message": "插入布局失败: [错误描述]",
               "data": {"original_slide_count": 10},
               "output_path": null
             }
    """
    logger.info(f"Executing add_slide_with_layout for: {prs_path}, layout: {layout_name}")
    result_dict = insert_layout(prs_path, layout_name, output_path, slide_title)
    return json.dumps(result_dict, ensure_ascii=False, indent=2)

@mcp.tool(description="Clears content from placeholders in specified slides of a presentation and returns a JSON string of the result.")
def clear_placeholders_from_slides(prs_path: str, output_path: Optional[str] = None, slide_indices: Optional[List[int]] = None) -> str:
    """
    清空指定PowerPoint演示文稿中特定或所有幻灯片内占位符的文本内容。

    此函数会遍历指定（或全部）幻灯片上的所有占位符，并尝试清除其文本内容。
    操作会保留占位符本身结构，仅移除文本。图片、表格等非文本内容不受影响。
    函数返回一个JSON字符串，包含操作结果，如处理的幻灯片数量、清空的占位符总数等。

    Args:
        prs_path (str): 源PPTX文件的绝对或相对路径。
        output_path (Optional[str], optional): 修改后PPTX的输出文件路径。
                                            如果为None，则会在标准输出目录下自动生成一个文件名，
                                            格式通常为 `[原文件名]_content_cleared_[时间戳].pptx`。
                                            默认为 None。
        slide_indices (Optional[List[int]], optional): 一个包含幻灯片索引（0-based）的列表，
                                                      指定要处理哪些幻灯片。
                                                      如果为None，则处理演示文稿中的所有幻灯片。
                                                      默认为 None。

    Returns:
        str: 一个JSON格式的字符串，包含了操作结果。
             成功时结构示例:
             {
               "status": "success",
               "message": "Successfully cleared content from 3 placeholder(s) in 2 slide(s).",
               "data": {
                 "slides_processed_count": 2, // 实际处理并有内容被清空的幻灯片数量
                 "placeholders_cleared_total": 3, // 所有被清空内容的占位符总数
                 "slides_targetted_count": 2, // 目标处理的幻灯片数量（基于slide_indices或总数）
                 "processed_slides_details": [ // 每个被处理幻灯片的详情
                   {
                     "slide_number": 1, // 幻灯片页码 (1-based)
                     "cleared_count_on_slide": 1, // 该幻灯片上被清空的占位符数量
                     "placeholders_status": [ // 该幻灯片上各占位符的处理状态
                       {"access_id": 0, "type": "标题 (Title)", "cleared": true, "reason": "text_frame cleared"},
                       {"access_id": 1, "type": "正文/内容 (Body)", "cleared": false, "reason": "no text content or not clearable type"}
                     ]
                   }
                 ]
               },
               "output_path": "path/to/output_content_cleared_timestamp.pptx"
             }
             失败时（例如文件处理错误）结构示例:
             {
               "status": "error",
               "message": "清空占位符内容失败: [错误描述]",
               "data": {"placeholders_cleared_total": 0, "slides_processed_count": 0},
               "output_path": null
             }
    """
    logger.info(f"Executing clear_placeholders_from_slides for: {prs_path}, slide_indices: {slide_indices}")
    result_dict = clear_placeholder_content(prs_path, output_path, slide_indices)
    return json.dumps(result_dict, ensure_ascii=False, indent=2)

@mcp.tool(description="Assigns content to a specific placeholder on a specific slide and returns a JSON string of the result.")
def set_placeholder_value(prs_path: str, slide_idx: int, placeholder_id: int, content_to_set: str, output_path: Optional[str] = None) -> str:
    """
    给指定PowerPoint演示文稿中特定幻灯片的特定占位符赋予文本内容。

    此函数通过幻灯片索引 (0-based) 和占位符的访问ID (placeholder_format.idx)
    来定位目标占位符，并将其文本内容设置为用户提供的 `content_to_set`。
    主要适用于文本类型的占位符（如标题、正文、副标题等）。
    操作结果以JSON字符串形式返回。

    Args:
        prs_path (str): 源PPTX文件的绝对或相对路径。
        slide_idx (int): 要修改的幻灯片的索引 (0-based)。
        placeholder_id (int): 要赋值的目标占位符的访问ID。这个ID即为占位符在其母版布局中定义的唯一 `idx` 值 
                              (即 `placeholder_format.idx` 属性值)。
                              **强烈建议通过先调用 `analyze_presentation_layouts` 工具来获取指定幻灯片上确切可用的占位符
                              及其对应的 `access_id`**（在 `analyze_presentation_layouts` 返回结果的 `data.masters_details.layouts.placeholders.access_id` 
                              或 `data.slides_details.placeholders_on_slide.access_id` 路径下可以找到，尽管后者需要先通过分析工具获取幻灯片上实际有哪些占位符），以确保操作的准确性。
        content_to_set (str): 要赋给占位符的文本内容。
        output_path (Optional[str], optional): 修改后PPTX的输出文件路径。
                                            如果为None，则会在标准输出目录下自动生成一个文件名，
                                            格式通常为 `[原文件名]_assigned_S[slide_idx]_P[placeholder_id]_[时间戳].pptx`。
                                            默认为 None。

    Returns:
        str: 一个JSON格式的字符串，包含了操作结果。
             成功时结构示例:
             {
               "status": "success",
               "message": "Successfully assigned content to placeholder ID 0 on slide 1.",
               "data": {
                 "slide_index": 0, // 目标幻灯片索引 (0-based)
                 "placeholder_access_id": 0, // 目标占位符访问ID
                 "content_assigned_length": 15, // 赋值内容的长度
                 "assignment_method": "text_frame" // 使用的赋值方法
               },
               "output_path": "path/to/output_assigned_S0_P0_timestamp.pptx"
             }
             失败时（例如索引超范围、占位符未找到、占位符不支持文本赋值）结构示例:
             {
               "status": "error",
               "message": "在幻灯片 0 上未找到访问ID为 99 的占位符。",
               "data": {"available_placeholders": [{"access_id": 0, "type": "标题 (Title)"}, ...]}, // 可能包含可用占位符信息
               "output_path": null
             }
             其他失败情况（例如文件读写错误）:
             {
               "status": "error",
               "message": "给占位符赋值时发生严重错误: [错误描述]",
               "data": null,
               "output_path": null
             }
    """
    logger.info(f"Executing set_placeholder_value for: {prs_path}, slide: {slide_idx}, placeholder: {placeholder_id}")
    result_dict = assign_placeholder_content(prs_path, slide_idx, placeholder_id, content_to_set, output_path)
    return json.dumps(result_dict, ensure_ascii=False, indent=2)

# 启动服务器
if __name__ == "__main__":
    # 确保必要的目录存在
    tmp_dir = get_tmp_dir()
    output_dir = get_output_dir()

    mcp.run(transport='stdio')