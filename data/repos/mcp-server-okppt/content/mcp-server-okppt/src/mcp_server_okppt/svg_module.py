import zipfile
import os
import uuid
import shutil
from lxml import etree
from reportlab.graphics import renderPM
from svglib.svglib import svg2rlg
from pptx.util import Inches, Pt, Cm, Emu
from typing import Optional, Union, Tuple, List
import traceback
import sys
from io import StringIO
import datetime

# 定义命名空间
ns = {
    'p': "http://schemas.openxmlformats.org/presentationml/2006/main",
    'a': "http://schemas.openxmlformats.org/drawingml/2006/main",
    'r': "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    'asvg': "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
}

# 路径辅助函数
def get_base_dir():
    """获取基础目录（服务器目录的父目录）"""
    current_dir = os.path.dirname(os.path.abspath(__file__))
    return os.path.dirname(current_dir)

def get_tmp_dir():
    """获取临时文件目录，如果不存在则创建"""
    tmp_dir = os.path.join(get_base_dir(), "tmp")
    os.makedirs(tmp_dir, exist_ok=True)
    return tmp_dir

def get_output_dir():
    """获取输出文件目录，如果不存在则创建"""
    output_dir = os.path.join(get_base_dir(), "output")
    os.makedirs(output_dir, exist_ok=True)
    return output_dir

def create_temp_dir():
    """创建唯一的临时目录并返回路径"""
    temp_dir = os.path.join(get_tmp_dir(), f"pptx_{uuid.uuid4().hex}")
    os.makedirs(temp_dir, exist_ok=True)
    return temp_dir

# 添加一个路径规范化函数
def normalize_path(path: str) -> str:
    """
    规范化路径格式，处理不同的路径表示方法，
    包括正斜杠、反斜杠、多重斜杠等情况。
    
    Args:
        path: 需要规范化的路径

    Returns:
        规范化后的路径
    """
    if not path:
        return path
    
    # 标准化路径，处理多重斜杠和混合斜杠的情况
    normalized = os.path.normpath(path)
    
    # 如果路径是绝对路径，转换为绝对路径
    if os.path.isabs(normalized):
        return normalized
    else:
        # 相对路径保持不变
        return normalized

# 添加一个创建SVG文件的函数
def create_svg_file(svg_path: str, width: int = 100, height: int = 100, text: str = "自动生成的SVG") -> bool:
    """
    创建一个简单的SVG文件。

    Args:
        svg_path: 要创建的SVG文件路径
        width: SVG宽度（像素）
        height: SVG高度（像素）
        text: 要在SVG中显示的文本

    Returns:
        bool: 如果成功创建则返回True，否则返回False
    """
    try:
        # 规范化路径
        svg_path = normalize_path(svg_path)
        
        # 获取文件目录并确保存在
        svg_dir = os.path.dirname(svg_path)
        if svg_dir and not os.path.exists(svg_dir):
            os.makedirs(svg_dir, exist_ok=True)
            print(f"已创建SVG目录: {svg_dir}")
            
        # 创建一个简单的SVG
        svg_content = (
            f'<svg width="{width}" height="{height}" xmlns="http://www.w3.org/2000/svg">\n'
            f'  <rect width="100%" height="100%" fill="#f0f0f0"/>\n'
            f'  <text x="{width/2}" y="{height/2}" font-size="14" text-anchor="middle" fill="black">{text}</text>\n'
            f'</svg>'
        )
        
        with open(svg_path, "w", encoding="utf-8") as f:
            f.write(svg_content)
        
        print(f"成功创建SVG文件: {svg_path}")
        return True
    except Exception as e:
        print(f"创建SVG文件时出错: {e}")
        traceback.print_exc()
        return False

def save_svg_code_to_file(
    svg_code: str,
    output_path: str = "",# 空字符串表示自动创建，否则使用绝对路径
    create_dirs: bool = True
) -> Tuple[bool, str, str]:
    """
    将SVG代码保存为SVG文件。

    Args:
        svg_code: SVG代码内容
        output_path: 输出文件路径，如果未指定，则生成一个带有时间戳的文件名
        create_dirs: 是否创建不存在的目录

    Returns:
        Tuple[bool, str, str]: (成功标志, 绝对路径, 错误消息)
    """
    try:
        # 如果未提供输出路径，则生成一个带有时间戳的文件名
        if not output_path or output_path == "":
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = os.path.join(get_output_dir(), f"svg_{timestamp}.svg")
        
        # 确保路径是绝对路径
        if not os.path.isabs(output_path):
            output_path = os.path.abspath(output_path)
        
        # 获取文件目录并确保存在
        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            if create_dirs:
                os.makedirs(output_dir, exist_ok=True)
                print(f"已创建目录: {output_dir}")
            else:
                return False, "", f"目录 {output_dir} 不存在"
        
        # 保存SVG代码到文件
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(svg_code)
        
        print(f"成功保存SVG代码到文件: {output_path}")
        return True, output_path, ""
    except Exception as e:
        error_message = f"保存SVG代码到文件时出错: {e}"
        print(error_message)
        traceback.print_exc()
        return False, "", error_message

# EMU 单位转换辅助函数
def to_emu(value: Union[Inches, Pt, Cm, Emu, int, float]) -> str:
    """将pptx单位或数值（假定为Pt）转换为EMU字符串"""
    # 先处理明确的 pptx.util 类型
    if isinstance(value, Inches):
        emu_val = value.emu
    elif isinstance(value, Cm):
        emu_val = value.emu
    elif isinstance(value, Pt):
        emu_val = value.emu
    elif isinstance(value, Emu):
        emu_val = value.emu
    elif isinstance(value, (int, float)):
        # 如果是纯数字，假设单位是 Pt
        emu_val = Pt(value).emu
    else:
        raise TypeError(f"Unsupported unit type for EMU conversion: {type(value)}")

    return str(int(emu_val)) # 确保返回整数的字符串形式

# 创建SVG文件的辅助函数
def create_svg_file(svg_path: str, width: int = 100, height: int = 100, text: str = "自动生成的SVG") -> bool:
    """
    创建一个简单的SVG文件。

    Args:
        svg_path: 要创建的SVG文件路径
        width: SVG宽度（像素）
        height: SVG高度（像素）
        text: 要在SVG中显示的文本

    Returns:
        bool: 如果成功创建则返回True，否则返回False
    """
    try:
        # 获取文件目录并确保存在
        svg_dir = os.path.dirname(svg_path)
        if svg_dir and not os.path.exists(svg_dir):
            os.makedirs(svg_dir, exist_ok=True)
            print(f"已创建SVG目录: {svg_dir}")
            
        # 创建一个简单的SVG
        svg_content = (
            f'<svg width="{width}" height="{height}" xmlns="http://www.w3.org/2000/svg">\n'
            f'  <rect width="100%" height="100%" fill="#f0f0f0"/>\n'
            f'  <text x="{width/2}" y="{height/2}" font-size="14" text-anchor="middle" fill="black">{text}</text>\n'
            f'</svg>'
        )
        
        with open(svg_path, "w", encoding="utf-8") as f:
            f.write(svg_content)
        
        print(f"成功创建SVG文件: {svg_path}")
        return True
    except Exception as e:
        print(f"创建SVG文件时出错: {e}")
        traceback.print_exc()
        return False


def insert_svg_to_pptx(
    pptx_path: str,
    svg_path: str,
    slide_number: int = 1,
    x: Optional[Union[Inches, Pt, Cm, Emu, int]] = None,
    y: Optional[Union[Inches, Pt, Cm, Emu, int]] = None,
    width: Optional[Union[Inches, Pt, Cm, Emu, int]] = None,
    height: Optional[Union[Inches, Pt, Cm, Emu, int]] = None,
    output_path: Optional[str] = None,
    create_if_not_exists: bool = True
) -> Union[bool, Tuple[bool, str]]:
    """
    将 SVG 图像插入到 PPTX 文件指定幻灯片的指定位置。
    **默认行为:** 如果不提供 `x`, `y`, `width`, `height` 参数，SVG 将被插入
    为幻灯片的全屏尺寸（位置 0,0，尺寸为幻灯片定义的宽度和高度）。

    **覆盖默认行为:**
    - 提供 `x` 和 `y` 来指定左上角位置 (使用 pptx.util 单位, 如 Inches(1), Pt(72))。
    - 提供 `width` 和/或 `height` 来指定尺寸 (使用 pptx.util 单位)。
    - 如果只提供了 `width` 而未提供 `height`，函数将尝试根据 SVG 的原始
      宽高比计算高度。如果无法获取宽高比，将使用幻灯片的默认高度。

    **实现方式:**
    此函数通过直接操作 PPTX 的内部 XML 文件来实现 SVG 插入。
    它会自动将 SVG 转换为 PNG 作为备用图像（用于旧版Office或不支持SVG的查看器），
    并将矢量 SVG 和 PNG 备用图都嵌入到 PPTX 文件中。

    Args:
        pptx_path: 原始 PPTX 文件的路径。如果文件不存在且create_if_not_exists为True，将自动创建。
        svg_path: 要插入的 SVG 文件的路径。
        slide_number: 要插入 SVG 的目标幻灯片编号 (从 1 开始)。
        x: 图片左上角的 X 坐标 (可选, 使用 pptx.util 单位)。默认为 0。
        y: 图片左上角的 Y 坐标 (可选, 使用 pptx.util 单位)。默认为 0。
        width: 图片的宽度 (可选, 使用 pptx.util 单位)。默认为幻灯片宽度。
        height: 图片的高度 (可选, 使用 pptx.util 单位)。默认为幻灯片高度。
                如果只提供了 width 而未提供 height，将尝试根据 SVG 原始宽高比计算。
        output_path: 输出 PPTX 文件的路径。如果为 None，将覆盖原始文件。
        create_if_not_exists: 如果为True且PPTX文件不存在，将自动创建一个新文件。

    Returns:
        Union[bool, Tuple[bool, str]]: 如果成功插入则返回 (True, "")，否则返回 (False, error_details)。
                               错误细节包含所有错误消息和堆栈跟踪。

    Raises:
        FileNotFoundError: 如果 pptx_path 不存在且 create_if_not_exists 为 False。
        FileNotFoundError: 如果 svg_path 无效。
        etree.XMLSyntaxError: 如果 PPTX 内部的 XML 文件损坏或格式错误。
        Exception: 其他潜在错误，如图库依赖问题或文件权限问题。

    Dependencies:
        - lxml: 用于 XML 处理。
        - reportlab 和 svglib: 用于将 SVG 转换为 PNG。
        - python-pptx: 主要用于方便的单位转换 (Inches, Pt 等)。
    """
    # 创建错误消息收集器
    error_log = []
    def log_error(message):
        """记录错误消息到错误日志列表"""
        error_log.append(message)
        print(message)  # 仍然打印到控制台

    # 创建临时目录 - 移到函数开始部分
    temp_dir = create_temp_dir()

    # 规范化并转换为绝对路径
    pptx_path = normalize_path(pptx_path)
    svg_path = normalize_path(svg_path)
    if output_path:
        output_path = normalize_path(output_path)
    
    # 转换为绝对路径
    if not os.path.isabs(pptx_path):
        pptx_path = os.path.abspath(pptx_path)
    
    if not os.path.isabs(svg_path):
        svg_path = os.path.abspath(svg_path)
    
    if output_path and not os.path.isabs(output_path):
        output_path = os.path.abspath(output_path)
    
    # 确保PPTX父目录存在
    pptx_dir = os.path.dirname(pptx_path)
    if not os.path.exists(pptx_dir):
        try:
            os.makedirs(pptx_dir, exist_ok=True)
            log_error(f"创建目录: {pptx_dir}")
        except Exception as e:
            error_msg = f"创建目录 {pptx_dir} 时出错: {e}"
            log_error(error_msg)
            log_error(traceback.format_exc())
            return False, "\n".join(error_log)
    
    # 如果有输出路径，也确保其父目录存在
    if output_path:
        output_dir = os.path.dirname(output_path)
        if not os.path.exists(output_dir):
            try:
                os.makedirs(output_dir, exist_ok=True)
                log_error(f"创建输出目录: {output_dir}")
            except Exception as e:
                error_msg = f"创建输出目录 {output_dir} 时出错: {e}"
                log_error(error_msg)
                log_error(traceback.format_exc())
                return False, "\n".join(error_log)
        
    # 输入验证并自动创建PPTX（如果需要）
    if not os.path.exists(pptx_path):
        if create_if_not_exists:
            try:
                from pptx import Presentation
                prs = Presentation()
                # 设置为16:9尺寸
                prs.slide_width = Inches(16)
                prs.slide_height = Inches(9)
                
                # 直接创建足够多的幻灯片
                for i in range(slide_number):
                    prs.slides.add_slide(prs.slide_layouts[6])  # 6是空白幻灯片
                
                prs.save(pptx_path)
                log_error(f"已重新创建PPTX文件: {pptx_path}，包含{slide_number}张幻灯片")
                import time
                time.sleep(0.5)
                # 再次尝试解压
                os.makedirs(temp_dir, exist_ok=True)  # 创建临时目录
                with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
                    zip_ref.extractall(temp_dir)
            except Exception as e:
                error_msg = f"创建PPTX文件时出错: {e}"
                log_error(error_msg)
                log_error(traceback.format_exc())
                return False, "\n".join(error_log)
        else:
            error_msg = f"PPTX file not found: {pptx_path}"
            log_error(error_msg)
            return False, "\n".join(error_log)
    
    # 检查并确保SVG文件的父目录存在
    svg_dir = os.path.dirname(svg_path)
    if not os.path.exists(svg_dir):
        try:
            os.makedirs(svg_dir, exist_ok=True)
            log_error(f"创建SVG目录: {svg_dir}")
        except Exception as e:
            error_msg = f"创建SVG目录 {svg_dir} 时出错: {e}"
            log_error(error_msg)
            log_error(traceback.format_exc())
            return False, "\n".join(error_log)
    
    # 检查PPTX文件是否至少有一张幻灯片，如果没有则添加一张
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        
        # 检查指定的slide_number是否超出现有幻灯片数量
        if slide_number > len(prs.slides):
            log_error(f"幻灯片编号{slide_number}超出现有幻灯片数量{len(prs.slides)}，将自动添加缺失的幻灯片")
            # 获取空白幻灯片布局
            blank_slide_layout = prs.slide_layouts[6]  # 6是空白幻灯片
            
            # 计算需要添加的幻灯片数量
            slides_to_add = slide_number - len(prs.slides)
            
            # 循环添加所需数量的幻灯片
            for _ in range(slides_to_add):
                prs.slides.add_slide(blank_slide_layout)
                log_error(f"已添加新的空白幻灯片，当前幻灯片数量: {len(prs.slides)}")
            
            # 保存文件
            prs.save(pptx_path)
            # 给文件写入一些时间
            import time
            time.sleep(0.5)
        elif len(prs.slides) == 0:
            log_error(f"PPTX文件 {pptx_path} 没有幻灯片，添加一张空白幻灯片")
            blank_slide_layout = prs.slide_layouts[6]  # 6是空白幻灯片
            slide = prs.slides.add_slide(blank_slide_layout)
            prs.save(pptx_path)
            # 给文件写入一些时间
            import time
            time.sleep(0.5)
    except Exception as e:
        error_msg = f"检查或添加幻灯片时出错: {e}"
        log_error(error_msg)
        # 如果是无效的PPTX文件，可能是因为文件损坏或不是PPTX格式
        if "File is not a zip file" in str(e) or "document not found" in str(e) or "Package not found" in str(e):
            log_error(f"PPTX文件 {pptx_path} 似乎不是有效的PowerPoint文件，尝试重新创建")
            try:
                # 确保目录存在
                os.makedirs(os.path.dirname(pptx_path), exist_ok=True)
                
                # 重新创建一个新的PPTX文件
                prs = Presentation()
                prs.slide_width = Inches(16)
                prs.slide_height = Inches(9)
                blank_slide_layout = prs.slide_layouts[6]
                
                # 直接创建足够多的幻灯片
                for i in range(slide_number):
                    prs.slides.add_slide(blank_slide_layout)
                    
                prs.save(pptx_path)
                log_error(f"已重新创建PPTX文件: {pptx_path}，包含{slide_number}张幻灯片")
            except Exception as e2:
                error_msg = f"重新创建PPTX文件时出错: {e2}"
                log_error(error_msg)
                log_error(traceback.format_exc())
                return False, "\n".join(error_log)
        else:
            # 其他类型的错误
            log_error(traceback.format_exc())
            return False, "\n".join(error_log)
    
    # 确保文件存在且大小不为0
    if not os.path.exists(pptx_path) or os.path.getsize(pptx_path) == 0:
        error_msg = f"错误：PPTX文件 {pptx_path} 不存在或大小为0"
        log_error(error_msg)
        return False, "\n".join(error_log)
    
    if not os.path.exists(svg_path):
        error_msg = f"SVG file not found: {svg_path}"
        log_error(error_msg)
        return False, "\n".join(error_log)

    # 确定输出路径和创建临时目录
    output_path = output_path or pptx_path
    # 临时目录已在函数开始部分创建，确保目录存在
    os.makedirs(temp_dir, exist_ok=True)

    default_width_emu = None
    default_height_emu = None

    try:
        # 解压 PPTX
        try:
            with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
                zip_ref.extractall(temp_dir)
        except zipfile.BadZipFile as e:
            error_msg = f"解压PPTX文件时出错: {e}"
            log_error(error_msg)
            log_error("尝试重新创建PPTX文件...")
            # 创建一个新的PPTX文件并再次尝试
            try:
                prs = Presentation()
                prs.slide_width = Inches(16)
                prs.slide_height = Inches(9)
                blank_slide_layout = prs.slide_layouts[6]
                
                # 直接创建足够多的幻灯片
                for i in range(slide_number):
                    prs.slides.add_slide(blank_slide_layout)
                    
                prs.save(pptx_path)
                log_error(f"已重新创建PPTX文件: {pptx_path}，包含{slide_number}张幻灯片")
                import time
                time.sleep(0.5)
                # 再次尝试解压
                with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
                    zip_ref.extractall(temp_dir)
            except Exception as e2:
                error_msg = f"重新创建和解压PPTX文件时出错: {e2}"
                log_error(error_msg)
                log_error(traceback.format_exc())
                return False, "\n".join(error_log)

        # --- 读取 presentation.xml 获取默认幻灯片尺寸 ---
        pres_path = os.path.join(temp_dir, "ppt", "presentation.xml")
        if os.path.exists(pres_path):
            try:
                pres_tree = etree.parse(pres_path)
                pres_root = pres_tree.getroot()
                sldSz = pres_root.find('p:sldSz', namespaces=ns)
                if sldSz is not None:
                    default_width_emu = sldSz.get('cx')
                    default_height_emu = sldSz.get('cy')
                    if not default_width_emu or not default_height_emu:
                        log_error("Warning: Could not read valid cx or cy from presentation.xml. Default size might be incorrect.")
                        default_width_emu = default_height_emu = None # Reset if invalid
                else:
                    log_error("Warning: <p:sldSz> element not found in presentation.xml. Cannot determine default slide size.")
            except etree.XMLSyntaxError as e:
                log_error(f"Warning: Could not parse presentation.xml: {e}. Cannot determine default slide size.")
        else:
            log_error("Warning: presentation.xml not found. Cannot determine default slide size.")
        
        # 如果无法获取默认尺寸，提供一个备用值（例如16:9宽屏的EMU值）
        if default_width_emu is None or default_height_emu is None:
            default_width_emu = "12192000" # 16 inches
            default_height_emu = "6858000" # 9 inches
            log_error(f"Warning: Using fallback default size: width={default_width_emu}, height={default_height_emu} EMU.")

        # --- 处理和转换图像 ---
        base_filename = f"image_{uuid.uuid4().hex}"
        media_dir = os.path.join(temp_dir, "ppt", "media")
        os.makedirs(media_dir, exist_ok=True)

        svg_filename = f"{base_filename}.svg"
        png_filename = f"{base_filename}.png"
        svg_target_path = os.path.join(media_dir, svg_filename)
        png_target_path = os.path.join(media_dir, png_filename)

        # 复制 SVG
        shutil.copy2(svg_path, svg_target_path)

        # 转换 SVG -> PNG (使用reportlab和svglib替代cairosvg)
        svg_width_px = None
        svg_height_px = None
        try:
            # 使用svglib和reportlab替代cairosvg
            drawing = svg2rlg(svg_path)
            renderPM.drawToFile(drawing, png_target_path, fmt="PNG")
            # 获取SVG尺寸信息
            svg_width_px = drawing.width
            svg_height_px = drawing.height
        except Exception as e:
            error_msg = f"Error converting SVG to PNG using reportlab/svglib: {e}"
            log_error(error_msg)
            log_error(traceback.format_exc())
            if os.path.exists(png_target_path): os.remove(png_target_path)
            if os.path.exists(svg_target_path): os.remove(svg_target_path)
            return False, "\n".join(error_log)

        # --- 计算最终尺寸 (EMU) ---
        # 位置
        x_emu = "0" if x is None else to_emu(x)
        y_emu = "0" if y is None else to_emu(y)

        # 宽度
        width_emu = default_width_emu if width is None else to_emu(width)

        # 高度
        if height is None:
            if width is None: # 完全默认，使用幻灯片高度
                height_emu = default_height_emu
            else: # 指定了宽度，尝试计算高度
                if svg_width_px and svg_height_px and svg_width_px > 0 and svg_height_px > 0:
                    aspect_ratio = svg_height_px / svg_width_px
                    height_emu_val = int(int(width_emu) * aspect_ratio)
                    height_emu = str(height_emu_val)
                    log_error(f"Info: Calculated height based on SVG aspect ratio: {height_emu} EMU")
                else:
                    log_error(f"Warning: Could not determine SVG aspect ratio. Using default height: {default_height_emu} EMU")
                    height_emu = default_height_emu
        else: # 用户指定了高度
            height_emu = to_emu(height)
            
        # --- 修改关系文件 (.rels) ---
        rels_path = os.path.join(temp_dir, "ppt", "slides", "_rels", f"slide{slide_number}.xml.rels")
        if not os.path.exists(rels_path):
             error_msg = f"Error: Relationship file not found for slide {slide_number}: {rels_path}"
             log_error(error_msg)
             if os.path.exists(png_target_path): os.remove(png_target_path)
             if os.path.exists(svg_target_path): os.remove(svg_target_path)
             return False, "\n".join(error_log)

        parser = etree.XMLParser(remove_blank_text=True)
        rels_tree = etree.parse(rels_path, parser)
        rels_root = rels_tree.getroot()

        # 查找最大的现有 rId
        max_rid_num = 0
        for rel in rels_root.findall('Relationship', namespaces=rels_root.nsmap):
            rid = rel.get('Id')
            if rid and rid.startswith('rId'):
                try:
                    num = int(rid[3:])
                    if num > max_rid_num:
                        max_rid_num = num
                except ValueError:
                    continue
        
        # 生成新的 rId
        rid_num_svg = max_rid_num + 1
        rid_num_png = max_rid_num + 2
        rId_svg = f"rId{rid_num_svg}"
        rId_png = f"rId{rid_num_png}"

        rel_ns = rels_root.nsmap.get(None)
        rel_tag = f"{{{rel_ns}}}Relationship" if rel_ns else "Relationship"

        # 创建 PNG 关系
        png_rel = etree.Element(
            rel_tag,
            Id=rId_png,
            Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
            Target=f"../media/{png_filename}"
        )
        # 创建 SVG 关系
        svg_rel = etree.Element(
            rel_tag,
            Id=rId_svg,
            Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image", # 仍然使用 image 类型
            Target=f"../media/{svg_filename}"
        )

        rels_root.append(png_rel)
        rels_root.append(svg_rel)

        rels_tree.write(rels_path, xml_declaration=True, encoding='UTF-8', standalone="yes")

        # --- 修改幻灯片文件 (slideX.xml) ---
        slide_path = os.path.join(temp_dir, "ppt", "slides", f"slide{slide_number}.xml")
        if not os.path.exists(slide_path):
             error_msg = f"Error: Slide file not found: {slide_path}"
             log_error(error_msg)
             if os.path.exists(png_target_path): os.remove(png_target_path)
             if os.path.exists(svg_target_path): os.remove(svg_target_path)
             return False, "\n".join(error_log)

        slide_tree = etree.parse(slide_path, parser)
        slide_root = slide_tree.getroot()
        
        # 查找或创建 spTree
        spTree = slide_root.find('.//p:spTree', namespaces=ns)
        if spTree is None:
            cSld = slide_root.find('p:cSld', namespaces=ns)
            if cSld is None:
                 error_msg = f"Error: Could not find <p:cSld> in slide {slide_number}."
                 log_error(error_msg)
                 if os.path.exists(png_target_path): os.remove(png_target_path)
                 if os.path.exists(svg_target_path): os.remove(svg_target_path)
                 return False, "\n".join(error_log)
            spTree = etree.SubElement(cSld, etree.QName(ns['p'], 'spTree'))
            max_nv_id = 0
            for elem in slide_root.xpath('.//p:cNvPr[@id]|.//p:cNvGrpSpPr[@id]', namespaces=ns):
                try:
                    nv_id = int(elem.get('id'))
                    if nv_id > max_nv_id:
                         max_nv_id = nv_id
                except (ValueError, TypeError):
                     continue
            group_shape_id = max_nv_id + 1
            nvGrpSpPr = etree.SubElement(spTree, etree.QName(ns['p'], 'nvGrpSpPr'))
            etree.SubElement(nvGrpSpPr, etree.QName(ns['p'], 'cNvPr'), id=str(group_shape_id), name="")
            etree.SubElement(nvGrpSpPr, etree.QName(ns['p'], 'cNvGrpSpPr'))
            etree.SubElement(nvGrpSpPr, etree.QName(ns['p'], 'nvPr'))
            grpSpPr = etree.SubElement(spTree, etree.QName(ns['p'], 'grpSpPr'))
            etree.SubElement(grpSpPr, etree.QName(ns['a'], 'xfrm'))

        max_shape_id = 0
        for elem in slide_root.xpath('.//p:cNvPr[@id]|.//p:cNvGrpSpPr[@id]|.//p:cNvSpPr[@id]', namespaces=ns):
            try:
                shape_id_val = int(elem.get('id'))
                if shape_id_val > max_shape_id:
                    max_shape_id = shape_id_val
            except (ValueError, TypeError):
                continue
        shape_id = max(max_shape_id + 1, 2) # 确保ID至少为2

        # 构建 p:pic 元素
        pic = etree.Element(etree.QName(ns['p'], 'pic'))

        # nvPicPr
        nvPicPr = etree.SubElement(pic, etree.QName(ns['p'], 'nvPicPr'))
        cNvPr = etree.SubElement(nvPicPr, etree.QName(ns['p'], 'cNvPr'), id=str(shape_id), name=f"Vector {shape_id}") # 更新名称

        # 添加 a16:creationId 扩展
        extLst_cNvPr = etree.SubElement(cNvPr, etree.QName(ns['a'], 'extLst'))
        ext_cNvPr = etree.SubElement(extLst_cNvPr, etree.QName(ns['a'], 'ext'), uri="{FF2B5EF4-FFF2-40B4-BE49-F238E27FC236}")
        a16_ns = "http://schemas.microsoft.com/office/drawing/2014/main"
        etree.register_namespace('a16', a16_ns)
        etree.SubElement(ext_cNvPr, etree.QName(a16_ns, 'creationId'), id="{" + str(uuid.uuid4()).upper() + "}")

        cNvPicPr = etree.SubElement(nvPicPr, etree.QName(ns['p'], 'cNvPicPr'))
        etree.SubElement(cNvPicPr, etree.QName(ns['a'], 'picLocks'), noChangeAspect="1")
        etree.SubElement(nvPicPr, etree.QName(ns['p'], 'nvPr'))

        # blipFill
        blipFill = etree.SubElement(pic, etree.QName(ns['p'], 'blipFill'))
        blip = etree.SubElement(blipFill, etree.QName(ns['a'], 'blip'), {etree.QName(ns['r'], 'embed'): rId_png})

        # 添加包含 svgBlip 的扩展列表
        extLst_blip = etree.SubElement(blip, etree.QName(ns['a'], 'extLst'))
        ext_blip_svg = etree.SubElement(extLst_blip, etree.QName(ns['a'], 'ext'), uri="{96DAC541-7B7A-43D3-8B79-37D633B846F1}")
        asvg_ns_uri = ns['asvg']
        etree.register_namespace('asvg', asvg_ns_uri)
        etree.SubElement(ext_blip_svg, etree.QName(asvg_ns_uri, 'svgBlip'), {etree.QName(ns['r'], 'embed'): rId_svg})

        stretch = etree.SubElement(blipFill, etree.QName(ns['a'], 'stretch'))
        etree.SubElement(stretch, etree.QName(ns['a'], 'fillRect'))

        # spPr (使用最终计算的 EMU 值)
        spPr = etree.SubElement(pic, etree.QName(ns['p'], 'spPr'))
        xfrm = etree.SubElement(spPr, etree.QName(ns['a'], 'xfrm'))
        etree.SubElement(xfrm, etree.QName(ns['a'], 'off'), x=x_emu, y=y_emu)
        etree.SubElement(xfrm, etree.QName(ns['a'], 'ext'), cx=width_emu, cy=height_emu)

        prstGeom = etree.SubElement(spPr, etree.QName(ns['a'], 'prstGeom'), prst="rect")
        etree.SubElement(prstGeom, etree.QName(ns['a'], 'avLst'))

        # 将 pic 添加到 spTree
        spTree.append(pic)

        # 写回 slide 文件
        slide_tree.write(slide_path, xml_declaration=True, encoding='UTF-8', standalone="yes")

        # --- 修改 [Content_Types].xml ---
        content_types_path = os.path.join(temp_dir, "[Content_Types].xml")
        if not os.path.exists(content_types_path):
            error_msg = f"Error: [Content_Types].xml not found at {content_types_path}"
            log_error(error_msg)
            if os.path.exists(png_target_path): os.remove(png_target_path)
            if os.path.exists(svg_target_path): os.remove(svg_target_path)
            return False, "\n".join(error_log)

        content_types_tree = etree.parse(content_types_path, parser)
        content_types_root = content_types_tree.getroot()
        ct_ns = content_types_root.nsmap.get(None) 
        ct_tag = f"{{{ct_ns}}}Default" if ct_ns else "Default"

        # 检查并添加 PNG 和 SVG Content Type (如果不存在)
        png_exists = any(default.get('Extension') == 'png' for default in content_types_root.findall(ct_tag, namespaces=content_types_root.nsmap))
        svg_exists = any(default.get('Extension') == 'svg' for default in content_types_root.findall(ct_tag, namespaces=content_types_root.nsmap))

        added_new_type = False
        if not png_exists:
            log_error("Info: Adding PNG Content Type to [Content_Types].xml")
            png_default = etree.Element(ct_tag, Extension="png", ContentType="image/png")
            content_types_root.append(png_default)
            added_new_type = True
        if not svg_exists:
            log_error("Info: Adding SVG Content Type to [Content_Types].xml")
            svg_default = etree.Element(ct_tag, Extension="svg", ContentType="image/svg+xml")
            content_types_root.append(svg_default)
            added_new_type = True

        if added_new_type:
            content_types_tree.write(content_types_path, xml_declaration=True, encoding='UTF-8', standalone="yes")

        # --- 重新打包 PPTX ---
        if os.path.exists(output_path):
            os.remove(output_path)

        with zipfile.ZipFile(output_path, 'w', compression=zipfile.ZIP_DEFLATED) as zip_out:
            for root, _, files in os.walk(temp_dir):
                for file in files:
                    file_path = os.path.join(root, file)
                    arcname = os.path.relpath(file_path, temp_dir)
                    zip_out.write(file_path, arcname)

        return True, ""

    except FileNotFoundError as e:
        error_msg = f"File not found error: {e}"
        log_error(error_msg)
        log_error(traceback.format_exc())
        return False, "\n".join(error_log)
    except etree.XMLSyntaxError as e:
        error_msg = f"XML parsing error: {e}"
        log_error(error_msg)
        log_error(traceback.format_exc())
        return False, "\n".join(error_log)
    except Exception as e:
        error_msg = f"An unexpected error occurred: {e}"
        log_error(error_msg)
        log_error(traceback.format_exc())
        return False, "\n".join(error_log)

    finally:
        # 清理临时目录
        try:
            shutil.rmtree(temp_dir, ignore_errors=True)
        except Exception as e:
            log_error(f"清理临时目录时出错: {e}")

def get_pptx_slide_count(pptx_path: str) -> Tuple[int, str]:
    """
    获取PPTX文件中的幻灯片数量。
    
    Args:
        pptx_path: PPTX文件路径
        
    Returns:
        Tuple[int, str]: 返回(幻灯片数量, 错误信息)的元组。
                       如果成功，错误信息为空字符串。
                       如果失败，幻灯片数量为0，错误信息包含详细错误。
    """
    error_message = ""
    
    try:
        # 规范化路径
        pptx_path = normalize_path(pptx_path)
        
        # 检查文件是否存在
        if not os.path.exists(pptx_path):
            return 0, f"文件不存在: {pptx_path}"
        
        # 使用python-pptx库打开文件并获取幻灯片数量
        from pptx import Presentation
        prs = Presentation(pptx_path)
        return len(prs.slides), ""
        
    except Exception as e:
        error_trace = traceback.format_exc()
        error_message = f"获取幻灯片数量时出错: {str(e)}\n{error_trace}"
        return 0, error_message

# --- 测试代码 ---
if __name__ == '__main__':
    import datetime
    
    # 测试用例：使用用户提供的 svg_test.pptx 和 media 目录下的 SVG
    input_pptx = "svg_test.pptx"
    
    # 模拟pptx_path为空的情况
    test_empty_path = False  # 将此设为True来测试空路径处理
    if test_empty_path:
        print("\n--- 测试空路径处理 ---")
        input_pptx = ""  # 模拟空路径
    
    # 检查pptx_path是否为空，如果为空则创建默认路径
    if not input_pptx or input_pptx.strip() == "":
        # 创建一个基于时间戳的默认文件名
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        input_pptx = f"presentation_{timestamp}.pptx"
        print(f"未提供PPTX路径，将使用默认路径: {input_pptx}")
    
    # 确保测试文件存在
    if not os.path.exists(input_pptx):
         print(f"Error: Test input file '{input_pptx}' not found in the workspace root.")
         # 可选：如果需要，可以自动创建
         from pptx import Presentation
         prs = Presentation()
         # 设置为16:9尺寸以便与默认备用尺寸匹配
         prs.slide_width = Inches(16)
         prs.slide_height = Inches(9)
         # 添加一张幻灯片，确保slide1.xml和相关关系文件存在
         blank_slide_layout = prs.slide_layouts[6]  # 6是空白幻灯片
         slide = prs.slides.add_slide(blank_slide_layout)
         prs.save(input_pptx)
         print(f"Created dummy file: {input_pptx} (16:9) with one slide")

    svg_to_insert = "image2.svg"  # 修改为当前目录下的SVG文件
    if not os.path.exists(svg_to_insert):
        print(f"Error: Test SVG file not found: {svg_to_insert}")
        # 可以在这里添加创建虚拟 SVG 的逻辑（如果需要）
        svg_content = (
            '<svg width="100" height="100" xmlns="http://www.w3.org/2000/svg">\n'
            '  <rect width="100%" height="100%" fill="orange"/>\n'
            '  <text x="50" y="60" font-size="15" text-anchor="middle" fill="black">Test SVG</text>\n'
            '</svg>'
        )
        try:
            with open(svg_to_insert, "w") as f:
                f.write(svg_content)
            print(f"Created dummy SVG file for testing: {svg_to_insert}")
        except Exception as e:
            print(f"Could not create dummy SVG: {e}")

    # --- 测试 1：指定尺寸和位置 --- (保持不变)
    output_pptx_specific = "svg_test_output_specific.pptx"
    print(f"\n--- Test 1: Inserting with specific size and position ---")
    if os.path.exists(input_pptx) and os.path.exists(svg_to_insert):
        success1, error_details1 = insert_svg_to_pptx(
            pptx_path=input_pptx,
            svg_path=svg_to_insert,
            slide_number=3,
            x=Inches(1),
            y=Inches(1),
            width=Inches(4),
            height=Inches(3),
            output_path=output_pptx_specific
        )
        if success1:
            print(f"SVG inserted with specific size/pos successfully into '{output_pptx_specific}'")
        else:
            print("Failed to insert SVG with specific size/pos.")
            print(error_details1)
    else:
        print("Skipping specific size test due to missing input files.")

    # --- 测试 2：默认全屏插入 --- (新增)
    output_pptx_fullscreen = "svg_test_output_fullscreen.pptx"
    print(f"\n--- Test 2: Inserting with default full-screen size ---")
    if os.path.exists(input_pptx) and os.path.exists(svg_to_insert):
        success2, error_details2 = insert_svg_to_pptx(
            pptx_path=input_pptx,
            svg_path=svg_to_insert,
            slide_number=1,
            # x, y, width, height 使用默认值 (None)
            output_path=output_pptx_fullscreen
        )
        if success2:
            print(f"SVG inserted with default full-screen successfully into '{output_pptx_fullscreen}'")
        else:
            print("Failed to insert SVG with default full-screen.")
            print(error_details2)
    else:
        print("Skipping full-screen test due to missing input files.")

    # --- 测试 3：只指定宽度，高度自动计算 --- (新增)
    output_pptx_autoheight = "svg_test_output_autoheight.pptx"
    print(f"\n--- Test 3: Inserting with specific width, auto height ---")
    if os.path.exists(input_pptx) and os.path.exists(svg_to_insert):
        success3, error_details3 = insert_svg_to_pptx(
            pptx_path=input_pptx,
            svg_path=svg_to_insert,
            slide_number=1,
            x=Inches(0.5),
            y=Inches(0.5),
            width=Inches(5), # 指定宽度
            # height 使用默认值 (None)
            output_path=output_pptx_autoheight
        )
        if success3:
            print(f"SVG inserted with auto height successfully into '{output_pptx_autoheight}'")
        else:
            print("Failed to insert SVG with auto height.")
            print(error_details3)
    else:
        print("Skipping auto height test due to missing input files.") 